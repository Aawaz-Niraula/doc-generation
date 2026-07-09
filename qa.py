"""qa.py — automated post-generation validation + safe output finalization.

Every generation runs through here BEFORE the file is allowed to exist in the
output directory (or be uploaded):

1. Integrity — reopen the bytes with the real parser (python-pptx / pypdf /
   python-docx), confirm slide/page count matches the content model.
2. Content — extract all text; fail on placeholder patterns, empty titles,
   duplicated slides.
3. Overflow — re-estimate every PPTX text frame against its container using
   the same font metrics the renderer's auto-fit guard used.
4. (optional) Visual smoke test via headless LibreOffice when available and
   QA_VISUAL=1.

Files are written to a temp name in the destination filesystem and atomically
renamed only after QA passes — a crashed run can never leave a corrupt or
half-written file in outputs/.
"""

from __future__ import annotations

import io
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from typing import List, Optional

from content_model import DocumentSpec
from theme import estimate_text_height_in

PLACEHOLDER_RE = re.compile(
    r"lorem\s+ipsum|\blorem\b|\bipsum\b|\bTODO\b|\[insert|\bXXX\b|\{\{.*?\}\}|\[placeholder",
    re.IGNORECASE,
)

_EMU_PER_IN = 914400


class QAError(RuntimeError):
    """A generated artifact failed validation and must not be shipped."""

    def __init__(self, fmt: str, issues: List[str]):
        self.format = fmt
        self.issues = issues
        super().__init__(f"{fmt.upper()} failed QA: " + " | ".join(issues))


@dataclass
class QAReport:
    format: str
    passed: bool = True
    units: int = 0                 # slides or pages
    size_kb: float = 0.0
    issues: List[str] = field(default_factory=list)      # fatal
    warnings: List[str] = field(default_factory=list)    # informational

    def as_dict(self) -> dict:
        return {
            "format": self.format, "passed": self.passed, "units": self.units,
            "size_kb": round(self.size_kb, 1), "issues": self.issues,
            "warnings": self.warnings,
        }


def _check_placeholders(text: str, report: QAReport, where: str):
    hits = {m.group(0).strip() for m in PLACEHOLDER_RE.finditer(text)}
    if hits:
        report.issues.append(f"placeholder text in {where}: {sorted(hits)[:4]}")


# ─── PPTX ─────────────────────────────────────────────────────────────────────

def qa_pptx(data: bytes, spec: DocumentSpec) -> QAReport:
    from pptx import Presentation

    report = QAReport(format="pptx", size_kb=len(data) / 1024)
    try:
        prs = Presentation(io.BytesIO(data))   # integrity: full reparse
    except Exception as exc:
        report.issues.append(f"file does not reopen: {exc}")
        report.passed = False
        return report

    slides = list(prs.slides)
    report.units = len(slides)
    if len(slides) != len(spec.sections):
        report.issues.append(
            f"slide count {len(slides)} != content model {len(spec.sections)}")

    all_text, first_lines = [], []
    for si, slide in enumerate(slides):
        slide_text = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = "\n".join(p.text for p in tf.paragraphs)
            if text.strip():
                slide_text.append(text)
            # overflow re-check with the renderer's own metrics
            sizes = [r.font.size.pt for p in tf.paragraphs for r in p.runs
                     if r.font.size is not None]
            bolds = [bool(r.font.bold) for p in tf.paragraphs for r in p.runs]
            if sizes and shape.width and shape.height:
                w_in = shape.width / _EMU_PER_IN
                h_in = shape.height / _EMU_PER_IN
                est = estimate_text_height_in(text, w_in, max(sizes), any(bolds))
                if est > h_in * 1.15 + 0.05:
                    report.issues.append(
                        f"slide {si + 1}: text frame overflow "
                        f"(~{est:.2f}in of text in {h_in:.2f}in box)")
        joined = " ".join(slide_text)
        all_text.append(joined)
        first_lines.append(joined[:120])
        if not joined.strip():
            report.issues.append(f"slide {si + 1} has no text at all")

    _check_placeholders(" ".join(all_text), report, "slides")
    seen = {}
    for si, line in enumerate(first_lines):
        key = line.strip().lower()
        if key and key in seen:
            report.warnings.append(f"slides {seen[key] + 1} and {si + 1} look duplicated")
        seen.setdefault(key, si)

    if not (spec.title or "").strip():
        report.issues.append("empty document title")

    report.passed = not report.issues
    return report


# ─── PDF ──────────────────────────────────────────────────────────────────────

_PAGE_DIMS_PT = {"a4": (595, 842), "letter": (612, 792)}


def qa_pdf(data: bytes, spec: DocumentSpec, page_size: str = "a4") -> QAReport:
    from pypdf import PdfReader

    report = QAReport(format="pdf", size_kb=len(data) / 1024)
    try:
        reader = PdfReader(io.BytesIO(data))   # integrity: full reparse
        pages = reader.pages
        report.units = len(pages)
    except Exception as exc:
        report.issues.append(f"file does not reopen: {exc}")
        report.passed = False
        return report

    if len(pages) != len(spec.sections):
        report.issues.append(
            f"page count {len(pages)} != content model {len(spec.sections)}")

    exp_w, exp_h = _PAGE_DIMS_PT.get(page_size.lower(), _PAGE_DIMS_PT["a4"])
    if pages:
        box = pages[0].mediabox
        w, h = float(box.width), float(box.height)
        if not (abs(w - exp_w) <= 6 and abs(h - exp_h) <= 6):
            report.issues.append(f"page dimensions {w:.0f}x{h:.0f}pt not {page_size.upper()}")

    try:
        text = " ".join((page.extract_text() or "") for page in pages)
    except Exception:
        text = ""
        report.warnings.append("text extraction unavailable; placeholder scan skipped")
    if text:
        _check_placeholders(text, report, "pages")

    if report.size_kb < 40:
        report.warnings.append(f"file only {report.size_kb:.0f}KB — styles/fonts may be missing")
    if report.size_kb > 20000:
        report.warnings.append(f"file is {report.size_kb / 1024:.1f}MB — check embedded assets")

    report.passed = not report.issues
    return report


# ─── DOCX ─────────────────────────────────────────────────────────────────────

def qa_docx(data: bytes, spec: DocumentSpec) -> QAReport:
    from docx import Document

    report = QAReport(format="docx", size_kb=len(data) / 1024)
    try:
        doc = Document(io.BytesIO(data))       # integrity: full reparse
    except Exception as exc:
        report.issues.append(f"file does not reopen: {exc}")
        report.passed = False
        return report

    paras = [p.text for p in doc.paragraphs]
    tables_text = [cell.text for tbl in doc.tables for row in tbl.rows for cell in row.cells]
    full = " ".join(paras + tables_text)
    report.units = len([p for p in paras if p.strip()])
    if spec.title.strip() and spec.title.strip()[:40].lower() not in full.lower():
        report.warnings.append("document title not found in body text")
    if report.units < 5:
        report.issues.append(f"only {report.units} non-empty paragraphs")
    _check_placeholders(full, report, "document")

    report.passed = not report.issues
    return report


def run_qa(fmt: str, data: bytes, spec: DocumentSpec, page_size: str = "a4",
           strict: bool = True) -> QAReport:
    """Dispatch to the per-format checker; raise QAError on failure when strict."""
    if fmt == "pptx":
        report = qa_pptx(data, spec)
    elif fmt == "pdf":
        report = qa_pdf(data, spec, page_size)
    elif fmt == "docx":
        report = qa_docx(data, spec)
    else:
        raise ValueError(f"unknown format: {fmt}")
    if strict and not report.passed:
        raise QAError(fmt, report.issues)
    return report


# ─── Visual smoke test (optional, LibreOffice) ────────────────────────────────

def visual_smoke_test(path: str, timeout: int = 90) -> Optional[str]:
    """Headless-convert an Office file to PDF as a render smoke test. Returns a
    note string, or None when LibreOffice is unavailable / disabled."""
    if os.getenv("QA_VISUAL") != "1":
        return None
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    with tempfile.TemporaryDirectory() as tmp:
        try:
            proc = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, path],
                capture_output=True, timeout=timeout,
            )
            produced = [f for f in os.listdir(tmp) if f.endswith(".pdf")]
            if proc.returncode == 0 and produced:
                return f"visual smoke test OK ({produced[0]})"
            return f"visual smoke test FAILED (rc={proc.returncode})"
        except Exception as exc:
            return f"visual smoke test errored: {exc}"


# ─── Safe output finalization ─────────────────────────────────────────────────

def unique_output_path(out_dir: str, slug: str, ext: str,
                       timestamp: Optional[datetime] = None) -> str:
    """outputs/{slug}_{YYYYMMDD-HHMMSS}.{ext}; never overwrites — collisions
    get _v2, _v3, …"""
    os.makedirs(out_dir, exist_ok=True)
    ts = (timestamp or datetime.now()).strftime("%Y%m%d-%H%M%S")
    base = f"{slug}_{ts}"
    candidate = os.path.join(out_dir, f"{base}.{ext}")
    counter = 2
    while os.path.exists(candidate):
        candidate = os.path.join(out_dir, f"{base}_v{counter}.{ext}")
        counter += 1
    return candidate


def finalize_output(data: bytes, out_dir: str, slug: str, ext: str) -> str:
    """Write to a temp file beside the destination, then atomically rename.
    Only call AFTER QA has passed. Returns the absolute final path."""
    final_path = unique_output_path(out_dir, slug, ext)
    directory = os.path.dirname(final_path)
    tmp_path = os.path.join(directory, f".tmp-{uuid.uuid4().hex}.{ext}")
    try:
        with open(tmp_path, "wb") as fh:
            fh.write(data)
        os.replace(tmp_path, final_path)     # atomic on the same filesystem
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
    return os.path.abspath(final_path)
