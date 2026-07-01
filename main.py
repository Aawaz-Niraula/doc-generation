from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Pt as PPTXPt, Inches as PPTXInches, Emu
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt as PPTXPt
from pptx.oxml import parse_xml as pptx_parse_xml
from pptx.oxml.ns import nsdecls as pptx_nsdecls, qn as pptx_qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from weasyprint import HTML, CSS
import httpx, io, json, os, re, uuid, math, traceback, random, asyncio, base64
from html import escape
from vercel.blob import AsyncBlobClient

app = FastAPI()

DEEPINFRA_KEY = os.getenv("DEEPINFRA_KEY")
VERCEL_BLOB_TOKEN = os.getenv("VERCEL_BLOB_TOKEN")


def _exception_detail(exc: Exception) -> str:
    trace_lines = traceback.format_exception(type(exc), exc, exc.__traceback__)
    tail = "".join(trace_lines[-8:]).strip()
    return f"{type(exc).__name__}: {str(exc) or 'Document generation failed'}\n{tail}"


@app.exception_handler(Exception)
async def unhandled_exception_handler(request, exc):
    return JSONResponse(
        status_code=500,
        content={"detail": _exception_detail(exc)},
    )


# ─── Health check ────────────────────────────────────────────────────────────

@app.get("/health")
async def health():
    return {"status": "ok"}


# ─── Utilities ───────────────────────────────────────────────────────────────

NUMBER_WORDS = {
    "one": 1, "two": 2, "three": 3, "four": 4, "five": 5, "six": 6,
    "seven": 7, "eight": 8, "nine": 9, "ten": 10, "eleven": 11, "twelve": 12,
    "thirteen": 13, "fourteen": 14, "fifteen": 15,
}

def requested_count(prompt: str, unit: str, default: int, minimum: int, maximum: int) -> int:
    normalized = prompt.lower()
    match = re.search(r"\b(\d{1,2})\s*(?:-|\s)?"+re.escape(unit)+r"s?\b", normalized)
    if match:
        return max(minimum, min(maximum, int(match.group(1))))
    for word, value in NUMBER_WORDS.items():
        if re.search(r"\b"+word+r"\s+(?:-|\s)?"+re.escape(unit)+r"s?\b", normalized):
            return max(minimum, min(maximum, value))
    return default


def _as_json(raw: str, label: str) -> dict:
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=502, detail=f"{label} returned invalid JSON") from exc


def _blob_value(blob, key: str, default=None):
    if isinstance(blob, dict):
        return blob.get(key, default)
    return getattr(blob, key, default)


def _generation_error(exc: Exception):
    if isinstance(exc, HTTPException):
        raise exc
    raise HTTPException(status_code=500, detail=_exception_detail(exc)) from exc


def _enforce_no_consecutive_layouts(items, type_key="type",
                                     interior_types=None,
                                     pin_first=None, pin_last=None):
    """Post-process AI output so no two consecutive items share the same layout."""
    if not items or len(items) < 3:
        return items
    if interior_types is None:
        interior_types = list({it.get(type_key, "") for it in items
                              if it.get(type_key, "") not in (pin_first, pin_last, "")})
    start = 1 if pin_first else 0
    end = len(items) - 1 if pin_last else len(items)
    for i in range(start, end):
        curr = items[i].get(type_key, "")
        prev = items[i - 1].get(type_key, "") if i > 0 else ""
        if curr == prev:
            nxt = items[i + 1].get(type_key, "") if i + 1 < len(items) else ""
            import random
            shuffled_alts = list(interior_types)
            random.shuffle(shuffled_alts)
            for alt in shuffled_alts:
                if alt != prev and alt != nxt:
                    items[i][type_key] = alt
                    break
    return items


async def _deepinfra_call(messages: list, temperature: float = 0.4) -> str:
    if not DEEPINFRA_KEY:
        raise HTTPException(status_code=500, detail="DEEPINFRA_KEY is not configured on the Render document service")

    async with httpx.AsyncClient(timeout=300) as client:
        res = await client.post(
            "https://api.deepinfra.com/v1/openai/chat/completions",
            headers={"Authorization": f"Bearer {DEEPINFRA_KEY}"},
            json={
                "model": "deepseek-ai/DeepSeek-V4-Flash",
                "temperature": temperature,
                "messages": messages,
            },
        )

    if res.status_code >= 400:
        raise HTTPException(
            status_code=502,
            detail=f"DeepInfra request failed ({res.status_code}): {res.text[:500]}",
        )

    try:
        raw = res.json()["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError, json.JSONDecodeError) as exc:
        raise HTTPException(status_code=502, detail="DeepInfra returned an invalid response") from exc

    return re.sub(r"```json|```", "", raw).strip()


# ─── Vercel Blob upload ───────────────────────────────────────────────────────

async def upload_to_vercel_blob(buffer: io.BytesIO, filename: str, content_type: str) -> dict:
    if not VERCEL_BLOB_TOKEN:
        raise HTTPException(status_code=500, detail="VERCEL_BLOB_TOKEN is not configured on the Render document service")
    buffer.seek(0)
    pathname = f"generated/{uuid.uuid4().hex}-{filename}"
    async with AsyncBlobClient(token=VERCEL_BLOB_TOKEN) as client:
        blob = await client.put(
            pathname, buffer.read(), access="public",
            content_type=content_type, add_random_suffix=False, overwrite=True,
        )
    url = _blob_value(blob, "url")
    if not url:
        raise HTTPException(status_code=502, detail="Vercel Blob upload returned no URL")

    download_url = _blob_value(blob, "download_url") or _blob_value(blob, "downloadUrl") or url
    return {
        "url": url,
        "downloadUrl": download_url,
        "pathname": _blob_value(blob, "pathname", pathname),
        "contentType": _blob_value(blob, "content_type") or _blob_value(blob, "contentType") or content_type,
    }


# ══════════════════════════════════════════════════════════════════════════════
# PDF GENERATION
# ══════════════════════════════════════════════════════════════════════════════

PALETTES = {
    "indigo":  {
        "primary":"#312E81","accent":"#6366F1","mid":"#818CF8","light":"#E0E7FF",
        "wash":"#F5F3FF","text":"#1E1B4B","dark":"#1E1B4B","soft":"#C7D2FE",
        "white":"#FFFFFF","body":"#374151","card2":"#4338CA","muted":"#A5B4FC",
    },
    "teal":    {
        "primary":"#134E4A","accent":"#0D9488","mid":"#2DD4BF","light":"#CCFBF1",
        "wash":"#F0FDFA","text":"#0F2F2E","dark":"#0F2F2E","soft":"#99F6E4",
        "white":"#FFFFFF","body":"#374151","card2":"#0F766E","muted":"#5EEAD4",
    },
    "crimson": {
        "primary":"#7F1D1D","accent":"#DC2626","mid":"#F87171","light":"#FEE2E2",
        "wash":"#FFF5F5","text":"#450A0A","dark":"#450A0A","soft":"#FCA5A5",
        "white":"#FFFFFF","body":"#374151","card2":"#991B1B","muted":"#FCA5A5",
    },
    "emerald": {
        "primary":"#064E3B","accent":"#059669","mid":"#34D399","light":"#D1FAE5",
        "wash":"#F0FDF4","text":"#022C22","dark":"#022C22","soft":"#6EE7B7",
        "white":"#FFFFFF","body":"#374151","card2":"#065F46","muted":"#6EE7B7",
    },
    "slate":   {
        "primary":"#1E293B","accent":"#475569","mid":"#94A3B8","light":"#E2E8F0",
        "wash":"#F8FAFC","text":"#0F172A","dark":"#0F172A","soft":"#CBD5E1",
        "white":"#FFFFFF","body":"#374151","card2":"#334155","muted":"#94A3B8",
    },
    "violet":  {
        "primary":"#4C1D95","accent":"#7C3AED","mid":"#A78BFA","light":"#EDE9FE",
        "wash":"#FAF5FF","text":"#2E1065","dark":"#2E1065","soft":"#C4B5FD",
        "white":"#FFFFFF","body":"#374151","card2":"#5B21B6","muted":"#C4B5FD",
    },
    "amber":   {
        "primary":"#78350F","accent":"#D97706","mid":"#FBBF24","light":"#FDE68A",
        "wash":"#FFFBEB","text":"#451A03","dark":"#451A03","soft":"#FCD34D",
        "white":"#FFFFFF","body":"#374151","card2":"#92400E","muted":"#FCD34D",
    },
    "rose":    {
        "primary":"#881337","accent":"#E11D48","mid":"#FB7185","light":"#FFE4E6",
        "wash":"#FFF1F2","text":"#4C0519","dark":"#4C0519","soft":"#FDA4AF",
        "white":"#FFFFFF","body":"#374151","card2":"#9F1239","muted":"#FDA4AF",
    },
    "navy":    {
        "primary":"#0F2D5E","accent":"#1D6FA4","mid":"#5BA8D4","light":"#D6EAF8",
        "wash":"#EBF5FB","text":"#0A1F3D","dark":"#0A1F3D","soft":"#AED6F1",
        "white":"#FFFFFF","body":"#374151","card2":"#154F7A","muted":"#7EC8E3",
    },
    "monochrome": {
        "primary":"#000000","accent":"#888888","mid":"#555555","light":"#E5E5E5",
        "wash":"#F5F5F5","text":"#111111","dark":"#000000","soft":"#CCCCCC",
        "white":"#FFFFFF","body":"#333333","card2":"#222222","muted":"#999999",
    },
}

PDF_PAGE_SEQUENCE = [
    "cover", "context", "deep_dive_1", "deep_dive_2", "deep_dive_3",
    "comparative", "case_study", "future", "takeaways", "closing",
]

PPTX_SLIDE_SEQUENCE = [
    "title", "context", "deep_dive_1", "deep_dive_2", "deep_dive_3",
    "comparative", "case_study", "future", "takeaways", "closing",
]

PDF_TYPE_TO_LAYOUT = {
    "cover": "cover", "context": "feature", "deep_dive_1": "editorial",
    "deep_dive_2": "split", "deep_dive_3": "stats", "comparative": "timeline",
    "case_study": "quote", "future": "manifesto", "takeaways": "grid", "closing": "closing",
    "split": "split", "feature": "feature", "grid": "grid", "quote": "quote",
    "timeline": "timeline", "stats": "stats", "editorial": "editorial",
    "manifesto": "manifesto",
}

PPTX_TYPE_TO_LAYOUT = {
    "title": "title", "cover": "title", "context": "stat", "deep_dive_1": "two_col",
    "deep_dive_2": "image_text", "deep_dive_3": "stat", "comparative": "two_col",
    "case_study": "quote", "future": "image_text", "takeaways": "stat", "closing": "closing",
    "bullets": "bullets", "two_col": "two_col", "stat": "stat", "quote": "quote",
    "image_text": "image_text",
}

AESTHETIC_DIRECTIONS = [
    "editorial_magazine", "scientific_brutalism", "cinematic_dark", "bauhaus_data",
    "japanese_minimalism", "retro_futurism", "nature_organic", "neon_cyberpunk",
    "luxury_editorial", "abstract_expressionism",
]


def _creative_director_core(doc_kind: str, unit_count: int, unit_label: str) -> str:
    aesthetics = ", ".join(AESTHETIC_DIRECTIONS)
    return f"""You are a world-class creative director, data visualization expert, and document strategist.
Return ONLY valid JSON — no markdown fences, no commentary.

CORE DIRECTIVE:
- Do NOT reuse generic templates. Pick ONE aesthetic direction from: {aesthetics}
  (or invent a better one). State which direction you chose and why it fits the topic.
- Commit fully to that aesthetic across all {unit_count} {unit_label}s.
- Use REAL topic-specific data: authentic statistics, dates, names, figures. No lorem ipsum.

DESIGN SYSTEM (define before content — lock it in):
- colors: exactly 5 named hex values — dominant, secondary, supporting_1, supporting_2, accent (accent used sparingly)
- typography: 2–3 Google Fonts with exact usage (h1_font, h2_font, body_font, data_font) plus weights/sizes
- layout_philosophy: one sentence describing grid rotation strategy
- signature_element: ONE visual motif repeated for cohesion (shape, line, texture, or chart style)

FORBIDDEN:
- Generic blue gradient backgrounds on every page
- Identical layout structure repeated
- Bullet lists as primary content format
- Default stock chart styling
- Centered title + centered body used more than once
- Placeholder or generic content

EVERY {unit_label.upper()} MUST include at least ONE of:
- chart_data (real numbers), hero_stat (single powerful stat), data_table (styled table),
  or shape_motif (describe SVG/geometric composition reinforcing the topic)

VISUAL STANDARD:
- Prefer data-native visuals, editorial geometry, and topic-specific diagrams over fake-looking AI photos.
- Use image_keyword only when a real-world visual materially helps the topic; never request fake people,
  fake screenshots, fake logos, fake UI, or readable text inside an image.
- Visuals must look intentionally art-directed, not like a generic AI template.

CONTENT STRUCTURE for {unit_count} {unit_label}s (adapt if count differs — preserve order):
1. COVER — cinematic, aesthetic fully declared
2. CONTEXT/OVERVIEW — stakes via data viz, not bullets
3. DEEP DIVE 1 — most complex layout, layered data
4. DEEP DIVE 2 — visual storytelling
5. DEEP DIVE 3 — quantitative, data-forward
6. COMPARATIVE or TIMELINE — side-by-side or chronological
7. CASE STUDY — pull quote + supporting data, editorial
8. FUTURE — forward-looking, abstract shapes
9. TAKEAWAYS — visual grid/cards, NOT bullet list
10. CLOSING — as impactful as cover, signature element"""


def _resolve_design_tokens(structure: dict) -> dict:
    """Merge AI design_system colors with legacy palette fallback."""
    ds = structure.get("design_system") or {}
    colors = ds.get("colors") or {}
    typo = ds.get("typography") or {}
    pname = str(structure.get("palette", "indigo")).lower()
    base = PALETTES.get(pname, PALETTES["indigo"])

    def _hex(key: str, fallback: str) -> str:
        val = colors.get(key) or colors.get(key.replace("_", ""))
        if val and re.match(r"^#[0-9A-Fa-f]{6}$", str(val)):
            return str(val)
        return fallback

    dominant = _hex("dominant", base["primary"])
    secondary = _hex("secondary", base["card2"])
    sup1 = _hex("supporting_1", base["mid"])
    sup2 = _hex("supporting_2", base["light"])
    accent = _hex("accent", base["accent"])

    return {
        "primary": dominant,
        "accent": accent,
        "mid": sup1,
        "light": sup2,
        "wash": _hex("supporting_2", base["wash"]),
        "text": base["text"],
        "dark": secondary if secondary != dominant else base["dark"],
        "soft": sup1,
        "white": base["white"],
        "card2": secondary,
        "muted": sup2,
        "body": base["body"],
        "h1_font": typo.get("h1_font") or typo.get("display_font") or "Playfair Display",
        "h2_font": typo.get("h2_font") or typo.get("h1_font") or "DM Sans",
        "body_font": typo.get("body_font") or "Source Sans 3",
        "data_font": typo.get("data_font") or typo.get("body_font") or "JetBrains Mono",
        "signature_element": ds.get("signature_element") or "",
        "layout_philosophy": ds.get("layout_philosophy") or "",
    }


def _google_fonts_link(tokens: dict) -> str:
    families = []
    for key in ("h1_font", "h2_font", "body_font", "data_font"):
        name = str(tokens.get(key, "")).strip()
        if name and name not in families:
            families.append(name)
    if not families:
        return ""
    params = "&".join(f"family={f.replace(' ', '+')}:wght@300;400;500;600;700;800;900" for f in families)
    return f'<link href="https://fonts.googleapis.com/css2?{params}&display=swap" rel="stylesheet"/>'


# ─── Real font embedding (self-contained, WeasyPrint-safe) ────────────────────
#
# WeasyPrint can mis-parse the Google Fonts CSS endpoint, so instead we fetch the
# actual static TTF files from the Fontsource CDN and inline them as base64
# @font-face rules. This makes the AI-chosen typography *actually render* in the
# PDF (previously it silently fell back to Georgia/Arial, which is the main reason
# every document looked identical). Fully time-boxed with graceful fallback.

_TTF_MAGIC = (b"\x00\x01\x00\x00", b"true", b"OTTO", b"ttcf")
_FONTSOURCE = "https://cdn.jsdelivr.net/fontsource/fonts/{slug}@latest/latin-{w}-normal.ttf"
_FONT_WEIGHTS = {
    "h1_font":   (700, 800, 900),
    "h2_font":   (500, 700),
    "body_font": (400, 600, 700),
    "data_font": (500, 700),
}
# A few common families whose CDN slug differs from a naive lower-hyphen transform.
_FONT_SLUG_FIXUPS = {
    "source sans 3": "source-sans-3", "source sans pro": "source-sans-pro",
    "ibm plex sans": "ibm-plex-sans", "ibm plex mono": "ibm-plex-mono",
    "pt sans": "pt-sans", "pt serif": "pt-serif",
}


def _font_slug(name: str) -> str:
    key = " ".join(str(name).strip().lower().split())
    if key in _FONT_SLUG_FIXUPS:
        return _FONT_SLUG_FIXUPS[key]
    return re.sub(r"[^a-z0-9]+", "-", key).strip("-")


async def _fetch_one_face(client: httpx.AsyncClient, family: str, weight: int):
    try:
        r = await client.get(_FONTSOURCE.format(slug=_font_slug(family), w=weight))
        if r.status_code == 200 and r.content[:4] in _TTF_MAGIC and len(r.content) < 900_000:
            b64 = base64.b64encode(r.content).decode("ascii")
            return (weight, f"@font-face{{font-family:'{family}';font-style:normal;"
                            f"font-weight:{weight};font-display:swap;"
                            f"src:url(data:font/ttf;base64,{b64}) format('truetype');}}")
    except Exception:
        pass
    return None


async def _build_embedded_font_css(tokens: dict) -> str:
    """Fetch the resolved typography families as base64 @font-face rules. Returns a
    <style> block, or '' on any failure so PDF rendering never depends on it."""
    families = {}
    for key, weights in _FONT_WEIGHTS.items():
        fam = str(tokens.get(key, "")).strip()
        if fam and re.search(r"[A-Za-z]", fam):
            families.setdefault(fam, set()).update(weights)
    if not families:
        return ""
    jobs = [(fam, w) for fam, ws in families.items() for w in sorted(ws)]
    try:
        async with httpx.AsyncClient(timeout=6.0, follow_redirects=True) as client:
            results = await asyncio.gather(
                *[_fetch_one_face(client, fam, w) for fam, w in jobs],
                return_exceptions=True,
            )
    except Exception:
        return ""
    faces = [r[1] for r in results if isinstance(r, tuple)]
    if not faces:
        return ""
    return "<style>" + "".join(faces) + "</style>"


def _generate_data_table_html(table: dict, P: dict) -> str:
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers or not rows:
        return ""
    hdr = "".join(
        f'<th style="background:{P["primary"]};color:white;padding:3mm 4mm;'
        f'font-family:var(--font-data);font-size:7.5pt;text-align:left;">{escape(str(h))}</th>'
        for h in headers[:6]
    )
    body_rows = ""
    for ri, row in enumerate(rows[:8]):
        bg = P["wash"] if ri % 2 == 0 else P["white"]
        cells = "".join(
            f'<td style="padding:2.8mm 4mm;font-family:var(--font-body);font-size:8.5pt;'
            f'border-bottom:0.3mm solid {P["light"]};">{escape(str(c))}</td>'
            for c in (list(row) + [""] * len(headers))[:len(headers)]
        )
        body_rows += f'<tr style="background:{bg};">{cells}</tr>'
    return f"""<div class="data-table-wrap" style="margin:6mm 0;border-radius:4mm;overflow:hidden;
      border:0.4mm solid {P['light']};">
  <table style="width:100%;border-collapse:collapse;">{hdr and f'<thead><tr>{hdr}</tr></thead>' or ''}
  <tbody>{body_rows}</tbody></table></div>"""


def _generate_hero_stat_html(stat: dict, P: dict) -> str:
    if not stat:
        return ""
    val = escape(str(stat.get("value") or stat.get("number") or "—"))
    lbl = escape(str(stat.get("label") or ""))
    ctx = escape(str(stat.get("context") or stat.get("sub") or ""))
    return f"""<div class="hero-stat" style="margin:8mm 0;padding:10mm 12mm;
      background:linear-gradient(135deg,{P['primary']},{P['accent']});
      border-radius:6mm;color:white;text-align:center;">
  <div style="font-family:var(--font-display);font-size:52pt;font-weight:800;
    letter-spacing:-0.04em;line-height:1;">{val}</div>
  <div style="font-family:var(--font-body);font-size:11pt;font-weight:600;
    letter-spacing:0.12em;text-transform:uppercase;margin-top:4mm;opacity:0.9;">{lbl}</div>
  <div style="font-family:var(--font-body);font-size:9pt;margin-top:3mm;opacity:0.75;max-width:120mm;margin-inline:auto;">{ctx}</div>
</div>"""


async def call_ai_pdf(prompt: str, page_count: int) -> dict:
    core = _creative_director_core("page", page_count, "page")
    system = f"""{core}

Generate exactly {page_count} PDF pages. Each page = one physical A4 sheet.

PAGE TYPES (fixed narrative order — each maps to a distinct grid):
- "cover"       → page 1: cinematic cover; declare aesthetic direction
- "context"     → overview with data visualization establishing stakes
- "deep_dive_1" → richest layout: layered data, annotations, chart or table
- "deep_dive_2" → visual storytelling: split narrative + highlights
- "deep_dive_3" → quantitative: 3 hero stat cards + chart
- "comparative" → timeline or before/after: steps with dates/figures
- "case_study"  → pull quote + attribution + supporting metrics
- "future"      → forward-looking manifesto statements + abstract tone
- "takeaways"   → 4 visual tiles/cards — NOT bullet list
- "closing"     → impactful close on signature element

Field usage:
- tiles: [{{"label":"","value":""}}] for takeaways/grid pages (exactly 4)
- steps: [{{"step":"","desc":"","date":""}}] for comparative/timeline (3-5)
- stats: [{{"number":"","label":"","sub":""}}] for deep_dive_3 (exactly 3)
- takeaways: ["visual card text"] for closing (3-4)
- statements: ["Bold sentence."] for future/manifesto (3)
- chart_data: [["Label", numeric_value], ...] (3-7 items) — real data
- hero_stat: {{"value":"847M","label":"METRIC","context":"one sentence"}}
- data_table: {{"headers":["Col1","Col2"],"rows":[["a","b"]]}}
- image_keyword: photorealistic topic image (when no chart/table on that page)

JSON schema:
{{
  "title": "Document title",
  "subtitle": "Compelling subtitle",
  "author": "Prepared for [topic] — 2026",
  "aesthetic_direction": {{
    "name": "cinematic_dark",
    "label": "Human-readable direction name",
    "rationale": "Why this aesthetic fits the topic (2 sentences)"
  }},
  "design_system": {{
    "colors": {{
      "dominant": "#1A1A2E", "secondary": "#16213E",
      "supporting_1": "#0F3460", "supporting_2": "#E8E8E8", "accent": "#E94560"
    }},
    "typography": {{
      "h1_font": "Playfair Display", "h2_font": "DM Sans",
      "body_font": "Source Sans 3", "data_font": "JetBrains Mono"
    }},
    "layout_philosophy": "Rotate full-bleed, 60/40 split, 3-column, hero, L-shape, diagonal",
    "signature_element": "Glowing accent arc repeated on every page footer"
  }},
  "palette": "indigo",
  "pages": [
    {{
      "type": "cover|context|deep_dive_1|deep_dive_2|deep_dive_3|comparative|case_study|future|takeaways|closing",
      "eyebrow": "SECTION LABEL",
      "heading": "Page headline",
      "body": "Rich paragraph, 80+ words, specific expert content with real figures.",
      "highlights": ["Insight one", "Insight two", "Insight three"],
      "callout": "Memorable quotable sentence.",
      "attribution": "— Source (case_study only)",
      "tiles": [], "steps": [], "stats": [], "takeaways": [], "statements": [],
      "chart_data": [], "hero_stat": {{}}, "data_table": {{}}, "image_keyword": ""
    }}
  ]
}}

RULES:
- pages[0].type MUST be "cover"
- pages[{page_count - 1}].type MUST be "closing"
- Follow narrative order above; skip middle sections only if page_count < 10
- Every page: include chart_data OR hero_stat OR data_table OR image_keyword
- body: 80+ words, topic-specific, no placeholders
- pages array MUST have exactly {page_count} items"""

    raw = await _deepinfra_call([
        {"role": "system", "content": system},
        {"role": "user", "content": prompt},
    ], temperature=0.38)
    data = _as_json(raw, "PDF content generator")
    pages = data.get("pages", [])
    while len(pages) < page_count:
        i = len(pages)
        seq_idx = min(i, len(PDF_PAGE_SEQUENCE) - 1)
        lt = PDF_PAGE_SEQUENCE[seq_idx]
        pages.append({
            "type": lt, "eyebrow": f"SECTION {i+1}",
            "heading": f"Key Insight {i+1}",
            "body": "This section presents essential analysis and actionable perspective on the topic, offering structured insight to guide informed decision-making and strategic planning across all relevant dimensions.",
            "highlights": ["Structured insight", "Clear implication", "Practical takeaway"],
            "callout": "The strongest documents stay precise, structured, and visually deliberate.",
            "tiles": [], "steps": [], "stats": [], "takeaways": [], "statements": [],
            "hero_stat": {}, "data_table": {},
        })
    data["pages"] = pages[:page_count]
    if page_count >= 2:
        data["pages"][0]["type"] = "cover"
        data["pages"][-1]["type"] = "closing"
    return data


async def _fetch_single_image(client: httpx.AsyncClient, url: str, page_dict: dict):
    try:
        resp = await client.get(url)
        if resp.status_code == 200:
            import base64
            b64 = base64.b64encode(resp.content).decode("utf-8")
            # Pollinations returns JPEG by default
            page_dict["image_data"] = f"data:image/jpeg;base64,{b64}"
    except Exception:
        # Silently fail, WeasyPrint will fall back to SVG
        pass


async def _prefetch_images(structure: dict):
    import asyncio
    import urllib.parse
    # Increased timeout to 25 seconds because Pollinations AI generates images on the fly. 
    # Since these run concurrently, this adds at most 25s total to the generation time.
    async with httpx.AsyncClient(timeout=25.0) as client:
        tasks = []
        for page in structure.get("pages", []):
            kw = page.get("image_keyword")
            if kw and str(kw).strip():
                encoded_kw = urllib.parse.quote(str(kw).strip() + " documentary editorial photograph, natural lighting, no text, no logo, high quality")
                url = f"https://image.pollinations.ai/prompt/{encoded_kw}?width=800&height=400&nologo=true"
                tasks.append(_fetch_single_image(client, url, page))
        if tasks:
            await asyncio.gather(*tasks, return_exceptions=True)


# ─── SVG decorative helpers (WeasyPrint renders inline SVG natively) ─────────

def _svg_circles(c1: str, c2: str, c3: str) -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;overflow:hidden;pointer-events:none;" viewBox="0 0 595 842" preserveAspectRatio="none">'
        f'<circle cx="520" cy="80" r="200" fill="{c1}" opacity="0.10"/>'
        f'<circle cx="50" cy="780" r="160" fill="{c2}" opacity="0.09"/>'
        f'<circle cx="480" cy="680" r="110" fill="{c3}" opacity="0.07"/>'
        f'<circle cx="280" cy="400" r="350" fill="{c1}" opacity="0.025"/>'
        f'</svg>'
    )

def _svg_diagonal_band(color: str, x1=0, y1=580, x2=595, y2=500, x3=595, y3=842, x4=0, y4=842) -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">'
        f'<polygon points="{x1},{y1} {x2},{y2} {x3},{y3} {x4},{y4}" fill="{color}" opacity="0.9"/>'
        f'</svg>'
    )

def _svg_grid_dots(color: str, spacing: int = 24) -> str:
    dots = "".join(
        f'<circle cx="{x}" cy="{y}" r="1.2" fill="{color}" opacity="0.30"/>'
        for x in range(0, 600, spacing) for y in range(0, 850, spacing)
    )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">{dots}</svg>'
    )

def _svg_half_circle(color: str, side: str = "right") -> str:
    if side == "right":
        return (
            f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;'
            f'top:0;right:0;width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">'
            f'<circle cx="595" cy="421" r="280" fill="{color}" opacity="0.12"/>'
            f'<circle cx="595" cy="421" r="180" fill="{color}" opacity="0.08"/>'
            f'</svg>'
        )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;'
        f'top:0;left:0;width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">'
        f'<circle cx="0" cy="421" r="280" fill="{color}" opacity="0.12"/>'
        f'<circle cx="0" cy="421" r="180" fill="{color}" opacity="0.08"/>'
        f'</svg>'
    )

def _svg_cross_hatch(color: str) -> str:
    lines = "".join(
        f'<line x1="{x}" y1="0" x2="{x}" y2="842" stroke="{color}" stroke-width="0.4" opacity="0.20"/>'
        for x in range(0, 600, 36)
    ) + "".join(
        f'<line x1="0" y1="{y}" x2="595" y2="{y}" stroke="{color}" stroke-width="0.4" opacity="0.20"/>'
        for y in range(0, 850, 36)
    )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">{lines}</svg>'
    )

def _svg_arc_decoration(color: str) -> str:
    arcs = "".join(
        f'<circle cx="595" cy="0" r="{r}" fill="none" stroke="{color}" stroke-width="1.2" opacity="{round(0.14 - i*0.025, 3)}"/>'
        for i, r in enumerate([100, 160, 220, 290, 360])
    )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">{arcs}</svg>'
    )

def _svg_stripe_band(color: str, accent: str) -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">'
        f'<polygon points="0,620 220,842 0,842" fill="{color}" opacity="0.16"/>'
        f'<polygon points="0,700 130,842 0,842" fill="{accent}" opacity="0.20"/>'
        f'</svg>'
    )

def _svg_diagonal_lines(color: str) -> str:
    """Subtle diagonal line pattern."""
    lines = "".join(
        f'<line x1="{x-200}" y1="0" x2="{x}" y2="842" stroke="{color}" stroke-width="0.6" opacity="0.12"/>'
        for x in range(-100, 900, 32)
    )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">{lines}</svg>'
    )

def _svg_hexagons(color: str) -> str:
    """Scattered hexagon outlines for texture."""
    hexes = ""
    positions = [(80,120,28),(500,80,20),(520,720,36),(60,700,22),(290,200,16),(450,450,24),(150,500,18)]
    for cx,cy,r in positions:
        pts = " ".join(
            f"{round(cx + r * math.cos(math.radians(60*i-30)),1)},{round(cy + r * math.sin(math.radians(60*i-30)),1)}"
            for i in range(6)
        )
        hexes += f'<polygon points="{pts}" fill="none" stroke="{color}" stroke-width="0.8" opacity="0.18"/>'
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">{hexes}</svg>'
    )

def _svg_noise_overlay(color: str) -> str:
    """Fine noise texture via tiny dots."""
    import random
    rng = random.Random(42)
    dots = "".join(
        f'<circle cx="{rng.randint(0,595)}" cy="{rng.randint(0,842)}" r="0.7" fill="{color}" opacity="0.08"/>'
        for _ in range(280)
    )
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;pointer-events:none;" viewBox="0 0 595 842">{dots}</svg>'
    )

def _svg_cover_geometric(primary: str, accent: str, mid: str) -> str:
    """Bold geometric shapes for cover pages."""
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;overflow:hidden;pointer-events:none;" viewBox="0 0 595 842">'
        # Large rotated rectangle top-right
        f'<rect x="320" y="-80" width="380" height="380" rx="18" fill="{accent}" opacity="0.13" transform="rotate(18 510 110)"/>'
        # Medium circle bottom-left
        f'<circle cx="40" cy="800" r="160" fill="{mid}" opacity="0.10"/>'
        # Small accent circle mid-right
        f'<circle cx="560" cy="500" r="70" fill="{accent}" opacity="0.09"/>'
        # Thin ring decoration
        f'<circle cx="510" cy="110" r="145" fill="none" stroke="{accent}" stroke-width="1.5" opacity="0.18"/>'
        f'<circle cx="510" cy="110" r="200" fill="none" stroke="{mid}" stroke-width="0.8" opacity="0.12"/>'
        f'</svg>'
    )

def _svg_manifesto_bg(dark: str, accent: str) -> str:
    """Dramatic background for manifesto pages."""
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" style="position:absolute;inset:0;'
        f'width:210mm;height:297mm;overflow:hidden;pointer-events:none;" viewBox="0 0 595 842">'
        f'<rect x="0" y="0" width="595" height="842" fill="{dark}"/>'
        f'<rect x="-60" y="200" width="200" height="600" fill="{accent}" opacity="0.07" transform="rotate(-12 40 500)"/>'
        f'<circle cx="595" cy="0" r="320" fill="{accent}" opacity="0.06"/>'
        f'<circle cx="0" cy="842" r="240" fill="{accent}" opacity="0.05"/>'
        f'</svg>'
    )

# ─── Data-driven SVG Charts ───────────────────────────────────────────────────

def _generate_svg_bar_chart(data: list, width=500, height=160, bar_color="#C9A84C", bg="#1E3A5F", title="") -> str:
    if not data: return ""
    uid = uuid.uuid4().hex[:6]
    max_val = max((d[1] for d in data), default=1)
    if max_val == 0: max_val = 1
    bars_html = ""
    w = width - 80
    h = height - 40
    bar_w = min(40, int(w / len(data)) - 10)
    for i, (label, val) in enumerate(data):
        bx = 40 + i * (w / len(data)) + (w / len(data) - bar_w) / 2
        bh = (val / max_val) * h
        by = height - 25 - bh
        bars_html += f'<rect x="{bx}" y="{by}" width="{bar_w}" height="{bh}" rx="4" fill="url(#barGrad{uid})" filter="url(#dropShadow{uid})"/>'
        bars_html += f'<text x="{bx+bar_w/2}" y="{by-5}" fill="#FFFFFF" font-family="Helvetica" font-size="9" text-anchor="middle" font-weight="bold">{val}</text>'
        bars_html += f'<text x="{bx+bar_w/2}" y="{height-10}" fill="rgba(255,255,255,0.7)" font-family="Helvetica" font-size="8" text-anchor="middle">{escape(str(label)[:8])}</text>'
    
    grid = "".join(f'<line x1="30" y1="{height-25 - (i*h/4)}" x2="{width-20}" y2="{height-25 - (i*h/4)}" stroke="rgba(255,255,255,0.1)" stroke-width="1"/>' for i in range(1, 5))
    
    return f"""<div style="margin-bottom:20px;">
    <svg width="{width}" height="{height}" viewBox="0 0 {width} {height}" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="barGrad{uid}" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="{bar_color}" />
                <stop offset="100%" stop-color="{bar_color}" stop-opacity="0.4" />
            </linearGradient>
            <filter id="dropShadow{uid}" x="-20%" y="-20%" width="140%" height="140%">
                <feDropShadow dx="0" dy="4" stdDeviation="3" flood-color="#000000" flood-opacity="0.3"/>
            </filter>
        </defs>
        <rect width="{width}" height="{height}" rx="12" fill="{bg}"/>
        {f'<text x="15" y="18" fill="#FFFFFF" font-family="Helvetica" font-size="11" font-weight="bold">{escape(title)}</text>' if title else ''}
        {grid}
        <line x1="30" y1="{height-25}" x2="{width-20}" y2="{height-25}" stroke="rgba(255,255,255,0.4)" stroke-width="1.5"/>
        {bars_html}
    </svg></div>"""

def _generate_svg_donut(segments: list, size=160, thickness=35, title="") -> str:
    if not segments: return ""
    total = sum(d[1] for d in segments)
    if total == 0: return ""
    cx, cy, r = size/2, size/2, (size - thickness)/2
    paths = ""
    start_angle = -math.pi / 2
    for label, val, color in segments:
        angle = (val / total) * 2 * math.pi
        end_angle = start_angle + angle
        x1 = cx + r * math.cos(start_angle)
        y1 = cy + r * math.sin(start_angle)
        x2 = cx + r * math.cos(end_angle)
        y2 = cy + r * math.sin(end_angle)
        large_arc = 1 if angle > math.pi else 0
        d = f"M {x1} {y1} A {r} {r} 0 {large_arc} 1 {x2} {y2}"
        paths += f'<path d="{d}" fill="none" stroke="{color}" stroke-width="{thickness}" stroke-linecap="butt"/>'
        # Add a tiny gap via stroke-dasharray if desired, but butt cap + math usually is fine, or we inject small white lines
        start_angle = end_angle
    
    return f"""<div style="text-align:center; margin-bottom:15px;">
    <svg width="{size}" height="{size}" viewBox="0 0 {size} {size}" xmlns="http://www.w3.org/2000/svg">
        <circle cx="{cx}" cy="{cy}" r="{r}" fill="none" stroke="rgba(0,0,0,0.05)" stroke-width="{thickness}"/>
        {paths}
        {f'<text x="{cx}" y="{cy}" fill="#333" font-family="Helvetica" font-size="14" font-weight="bold" text-anchor="middle" dominant-baseline="middle">{escape(title)}</text>' if title else ''}
    </svg></div>"""

def _generate_svg_area_chart(data: list, width=500, height=130, line_color="#4A90D9", title="") -> str:
    if not data: return ""
    uid = uuid.uuid4().hex[:6]
    max_val = max((d[1] for d in data), default=1)
    if max_val == 0: max_val = 1
    w = width - 40
    h = height - 30
    pts = []
    circles = ""
    labels = ""
    for i, (label, val) in enumerate(data):
        px = 20 + i * (w / max(1, len(data)-1))
        py = height - 20 - (val / max_val) * h
        pts.append((px, py))
        circles += f'<circle cx="{px}" cy="{py}" r="4" fill="{line_color}" stroke="#FFF" stroke-width="1.5"/>'
        labels += f'<text x="{px}" y="{height-5}" fill="rgba(0,0,0,0.6)" font-family="Helvetica" font-size="8" text-anchor="middle">{escape(str(label)[:8])}</text>'
    
    path_d = "M " + " L ".join(f"{x},{y}" for x, y in pts)
    area_d = f"M {pts[0][0]},{height-20} L " + " L ".join(f"{x},{y}" for x, y in pts) + f" L {pts[-1][0]},{height-20} Z"

    return f"""<div style="margin-bottom:20px;">
    <svg width="{width}" height="{height}" viewBox="0 0 {width} {height}" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <linearGradient id="areaGrad{uid}" x1="0" y1="0" x2="0" y2="1">
                <stop offset="0%" stop-color="{line_color}" stop-opacity="0.4"/>
                <stop offset="100%" stop-color="{line_color}" stop-opacity="0.02"/>
            </linearGradient>
        </defs>
        <path d="{area_d}" fill="url(#areaGrad{uid})"/>
        <path d="{path_d}" fill="none" stroke="{line_color}" stroke-width="2.5" stroke-linecap="round"/>
        {circles}
        {labels}
        <line x1="20" y1="{height-20}" x2="{width-20}" y2="{height-20}" stroke="rgba(0,0,0,0.1)" stroke-width="1"/>
    </svg></div>"""

# ─── Footer / running elements ────────────────────────────────────────────────

def _footer(title: str, idx: int, total: int, accent: str, mid: str, dark_bg: bool = False) -> str:
    title_color = "rgba(255,255,255,0.5)" if dark_bg else mid
    rule_color  = "rgba(255,255,255,0.15)" if dark_bg else "var(--light)"
    return f"""<div class="footer" style="border-top-color:{rule_color};">
  <span class="footer-title" style="color:{title_color};">{title}</span>
  <span class="footer-pager" style="background:{accent};color:white;">{idx}&thinsp;/&thinsp;{total}</span>
</div>"""

def _eyebrow(text: str, accent: str, light: bool = False) -> str:
    color = "rgba(255,255,255,0.65)" if light else accent
    return f'<div class="eyebrow" style="color:{color};">{text}</div>'


def _clip_words(value: str, limit: int) -> str:
    words = str(value or "").split()
    if len(words) <= limit:
        return " ".join(words)
    return " ".join(words[:limit]).rstrip(".,;:") + "."


def _clip_chars(value: str, limit: int) -> str:
    text = " ".join(str(value or "").split())
    if len(text) <= limit:
        return text
    return text[:limit].rsplit(" ", 1)[0].rstrip(".,;:") + "."


def _safe_chart_data(chart_data: list, limit: int = 7) -> list:
    values = []
    if not isinstance(chart_data, list):
        return values
    for item in chart_data[:limit]:
        try:
            label, val = item[0], float(item[1])
            values.append((str(label), val))
        except (TypeError, ValueError, IndexError):
            continue
    return values


def _premium_chart_html(data: list, P: dict, title: str = "", big: bool = False) -> str:
    data = _safe_chart_data(data)
    if not data:
        return ""
    max_val = max((v for _, v in data), default=1) or 1
    lo = min((v for _, v in data), default=0)
    rows = []
    for i, (label, val) in enumerate(data):
        width = max(7, min(100, (val / max_val) * 100))
        peak = " peak" if val == max_val else ""
        rows.append(f"""
<div class="chart-row">
  <div class="chart-label">{escape(_clip_chars(label, 22))}</div>
  <div class="chart-track"><div class="chart-bar{peak}" style="width:{width:.1f}%"></div></div>
  <div class="chart-value">{escape(f'{val:g}')}</div>
</div>""")
    cls = "premium-chart big" if big else "premium-chart"
    delta = ""
    if len(data) >= 2 and data[0][1]:
        try:
            chg = (data[-1][1] - data[0][1]) / abs(data[0][1]) * 100
            arrow = "▲" if chg >= 0 else "▼"
            delta = f'<span class="chart-delta">{arrow} {abs(chg):.0f}%</span>'
        except Exception:
            delta = ""
    return f"""
<div class="{cls}">
  <div class="chart-head"><span class="chart-kicker">{escape(_clip_chars(title or 'Measured signal', 40))}</span>{delta}</div>
  <div class="chart-rows">{"".join(rows)}</div>
</div>"""


def _premium_stat_html(stat: dict, P: dict) -> str:
    """Large gradient hero-stat card (used as a page centerpiece)."""
    if not isinstance(stat, dict) or not (stat.get("value") or stat.get("number")):
        return ""
    return f"""
<div class="premium-stat">
  <div class="premium-stat-value">{escape(str(stat.get("value") or stat.get("number")))}</div>
  <div class="premium-stat-label">{escape(_clip_chars(stat.get("label") or "Key metric", 42))}</div>
  <p>{escape(_clip_words(stat.get("context") or stat.get("sub") or "", 26))}</p>
</div>"""


def _big_stat_card(stat: dict, P: dict, index: int = 0) -> str:
    """Tall stat tile with an oversized figure — for stat-grid pages."""
    val = escape(str(stat.get("value") or stat.get("number") or "—"))
    lbl = escape(_clip_chars(stat.get("label") or stat.get("value") or "Metric", 24))
    ctx = escape(_clip_words(stat.get("context") or stat.get("sub") or "", 18))
    return f"""
<div class="bigstat">
  <div class="bigstat-idx">{index+1:02d}</div>
  <div class="bigstat-val">{val}</div>
  <div class="bigstat-lbl">{lbl}</div>
  <div class="bigstat-ctx">{ctx}</div>
</div>"""


def _premium_table_html(table: dict, P: dict) -> str:
    if not isinstance(table, dict):
        return ""
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers or not rows:
        return ""
    hdr = "".join(f"<th>{escape(_clip_chars(h, 22))}</th>" for h in headers[:5])
    body = ""
    for row in rows[:7]:
        cells = "".join(
            f"<td>{escape(_clip_chars(c, 34))}</td>"
            for c in (list(row) + [""] * len(headers))[:len(headers[:5])]
        )
        body += f"<tr>{cells}</tr>"
    return f'<table class="premium-table"><thead><tr>{hdr}</tr></thead><tbody>{body}</tbody></table>'


def _geo_panel(P: dict, idx: int, label: str, dark: bool = True, tall: bool = False, variant: int = None) -> str:
    """Clean, art-directed geometric data-visual that fills a hero or side area
    (no blurry radial glows). Tall columns use crop-safe centered compositions;
    wide areas cycle through four distinct data-shaped motifs."""
    uid = f"gp{idx}"
    acc, mid, prim = P["accent"], P["mid"], P["primary"]
    ink = "rgba(255,255,255,.9)" if dark else prim
    faint = "rgba(255,255,255,.14)" if dark else "rgba(15,23,42,.10)"
    rng = random.Random(idx * 7 + 3)

    if tall:
        # Portrait aspect (matches the ~76×150mm side column) with a centered,
        # crop-safe motif so nothing important is clipped.
        W, H = 380, 750
        variant = (idx if variant is None else variant) % 2
        cx, cy = W / 2, H / 2
        if variant == 0:  # concentric arcs
            rings = "".join(
                f'<circle cx="{cx:.0f}" cy="{cy:.0f}" r="{r}" fill="none" stroke="{acc}" stroke-width="{w}" opacity="{o}" stroke-dasharray="{d}" stroke-linecap="round"/>'
                for r, w, o, d in [(160,"2",".16","0"),(128,"16",".26",f"{rng.randint(220,420)} 1200"),(96,"12",".42",f"{rng.randint(150,300)} 1200"),(64,"9",".62",f"{rng.randint(90,220)} 1200")])
            core = f'<circle cx="{cx:.0f}" cy="{cy:.0f}" r="30" fill="{acc}" opacity=".9"/><circle cx="{cx:.0f}" cy="{cy:.0f}" r="15" fill="{"#fff" if dark else prim}" opacity=".9"/>'
            ticks = "".join(f'<circle cx="{cx + 200*math.cos(math.radians(a)):.0f}" cy="{cy + 200*math.sin(math.radians(a)):.0f}" r="3" fill="{mid}" opacity=".5"/>' for a in range(0, 360, 30))
            art = ticks + rings + core
        else:  # stacked horizontal bars
            n = 7
            gap = (H - 160) / n
            bars = ""
            for i in range(n):
                bw = rng.randint(120, W - 90)
                y = 100 + i * gap
                c = acc if i % 3 == 0 else mid
                bars += f'<rect x="60" y="{y:.0f}" width="{bw}" height="{gap*0.5:.0f}" rx="6" fill="{c}" opacity="{0.45 + i*0.06:.2f}"/>'
            art = bars
    else:
        W, H = 600, 335
        variant = ((idx - 1) if variant is None else variant) % 4

    if not tall and variant == 0:  # ascending area + line
        n = 8
        ys = []
        y = H - 90
        for i in range(n):
            y = max(70, min(H - 70, y - rng.randint(-8, 46)))
            ys.append(y)
        xs = [60 + i * ((W - 120) / (n - 1)) for i in range(n)]
        line = "M " + " L ".join(f"{x:.0f},{y:.0f}" for x, y in zip(xs, ys))
        area = f"M {xs[0]:.0f},{H-60} L " + " L ".join(f"{x:.0f},{y:.0f}" for x, y in zip(xs, ys)) + f" L {xs[-1]:.0f},{H-60} Z"
        grid = "".join(f'<line x1="60" y1="{H-60-i*((H-130)/4):.0f}" x2="{W-40}" y2="{H-60-i*((H-130)/4):.0f}" stroke="{faint}" stroke-width="1"/>' for i in range(5))
        dots = "".join(f'<circle cx="{x:.0f}" cy="{y:.0f}" r="5" fill="{acc}" stroke="{ "#fff" if dark else prim}" stroke-width="2"/>' for x, y in zip(xs, ys))
        art = (f'<defs><linearGradient id="{uid}a" x1="0" y1="0" x2="0" y2="1">'
               f'<stop offset="0%" stop-color="{acc}" stop-opacity=".38"/>'
               f'<stop offset="100%" stop-color="{acc}" stop-opacity="0"/></linearGradient></defs>'
               f'{grid}<path d="{area}" fill="url(#{uid}a)"/>'
               f'<path d="{line}" fill="none" stroke="{acc}" stroke-width="3.5" stroke-linejoin="round" stroke-linecap="round"/>{dots}')
    elif not tall and variant == 1:  # column bars
        n = 7
        bw = (W - 130) / n
        bars = ""
        for i in range(n):
            bh = rng.randint(70, H - 130)
            x = 65 + i * bw
            y = H - 60 - bh
            c = acc if i % 3 == 0 else mid
            bars += f'<rect x="{x:.0f}" y="{y:.0f}" width="{bw*0.62:.0f}" height="{bh}" rx="5" fill="{c}" opacity="{0.55 + i*0.05:.2f}"/>'
        grid = "".join(f'<line x1="55" y1="{H-60-i*((H-130)/4):.0f}" x2="{W-40}" y2="{H-60-i*((H-130)/4):.0f}" stroke="{faint}" stroke-width="1"/>' for i in range(5))
        art = f'{grid}{bars}<line x1="55" y1="{H-60}" x2="{W-40}" y2="{H-60}" stroke="{faint}" stroke-width="1.5"/>'
    elif not tall and variant == 2:  # concentric arcs / radial
        cx, cy = W/2, H/2 + 6
        rings = "".join(
            f'<circle cx="{cx:.0f}" cy="{cy:.0f}" r="{r}" fill="none" stroke="{acc}" stroke-width="{w}" opacity="{o}" stroke-dasharray="{d}" stroke-linecap="round"/>'
            for r, w, o, d in [(150,"2",".18","0"),(120,"14",".28",f"{rng.randint(180,320)} 900"),(92,"11",".42",f"{rng.randint(120,240)} 900"),(64,"8",".6",f"{rng.randint(80,180)} 900")])
        core = f'<circle cx="{cx:.0f}" cy="{cy:.0f}" r="30" fill="{acc}" opacity=".9"/><circle cx="{cx:.0f}" cy="{cy:.0f}" r="15" fill="{"#fff" if dark else prim}" opacity=".85"/>'
        art = rings + core
    elif not tall:  # node network
        nodes = [(120,110,15),(260,70,11),(420,120,17),(510,210,12),(360,250,14),(180,250,10),(300,165,20),(470,60,9)]
        edges = [(0,6),(1,6),(2,6),(2,3),(3,4),(4,5),(5,0),(6,4),(1,2),(2,7)]
        lines = "".join(f'<line x1="{nodes[a][0]}" y1="{nodes[a][1]}" x2="{nodes[b][0]}" y2="{nodes[b][1]}" stroke="{faint}" stroke-width="1.4"/>' for a, b in edges)
        dots = "".join(f'<circle cx="{x}" cy="{y}" r="{r}" fill="{acc if i%2==0 else mid}" opacity="{0.5+i*0.05:.2f}"/><circle cx="{x}" cy="{y}" r="{max(3,r-6)}" fill="#fff" opacity=".55"/>' for i,(x,y,r) in enumerate(nodes))
        art = lines + dots

    bg = (f'<defs><linearGradient id="{uid}bg" x1="0" y1="0" x2="1" y2="1">'
          f'<stop offset="0%" stop-color="{prim}"/><stop offset="100%" stop-color="{P["card2"]}"/></linearGradient></defs>'
          f'<rect width="{W}" height="{H}" rx="20" fill="url(#{uid}bg)"/>') if dark else \
         f'<rect width="{W}" height="{H}" rx="20" fill="{P["wash"]}"/><rect width="{W}" height="{H}" rx="20" fill="none" stroke="{P["light"]}" stroke-width="1.5"/>'
    lab = (f'<rect x="34" y="34" width="30" height="7" rx="3.5" fill="{acc}"/>'
           f'<text x="34" y="{H-34}" fill="{ink}" font-family="monospace" font-size="16" font-weight="700" letter-spacing="2">{escape(_clip_chars(label, 34)).upper()}</text>')
    return (f'<div class="geo-panel"><svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
            f'preserveAspectRatio="xMidYMid slice">{bg}{art}{lab}</svg></div>')


def _hero_visual_html(page: dict, P: dict, idx: int, heading: str, dark: bool = False) -> str:
    """The dominant, page-filling centerpiece: chart > table > hero-stat > photo > geometric."""
    chart = _premium_chart_html(page.get("chart_data") or [], P, heading, big=True)
    if chart:
        return chart
    table = _premium_table_html(page.get("data_table") or {}, P)
    if table:
        return f'<div class="hero-table">{table}</div>'
    stat = _premium_stat_html(page.get("hero_stat") or {}, P)
    if stat:
        return stat
    image_data = str(page.get("image_data") or "").strip()
    if image_data:
        return f'<div class="photo-plate big"><img src="{image_data}" alt="" /><div class="photo-caption">{escape(_clip_chars(heading, 70))}</div></div>'
    return _geo_panel(P, idx, heading, dark=True, variant=page.get("_geo"))


def _side_visual_html(page: dict, P: dict, idx: int, heading: str, dark: bool = False) -> str:
    """A visual sized for a narrower side column (chart/stat/table/geometric)."""
    stat = _premium_stat_html(page.get("hero_stat") or {}, P)
    if stat:
        return stat
    table = _premium_table_html(page.get("data_table") or {}, P)
    if table:
        return f'<div class="hero-table">{table}</div>'
    chart = _premium_chart_html(page.get("chart_data") or [], P, heading, big=False)
    if chart:
        return chart
    image_data = str(page.get("image_data") or "").strip()
    if image_data:
        return f'<div class="photo-plate"><img src="{image_data}" alt="" /><div class="photo-caption">{escape(_clip_chars(heading, 60))}</div></div>'
    return _geo_panel(P, idx, heading, dark=True, tall=True, variant=page.get("_geo"))


def _insight_cards(highlights: list, P: dict, limit: int = 3) -> str:
    """A fixed row of equal-width insight cards (table-cell layout, never wraps)."""
    items = [h for h in (highlights or []) if str(h).strip()][:limit]
    if not items:
        return ""
    cells = "".join(
        f'<div class="cell"><div class="ins-card"><span class="ins-num">{i+1:02d}</span>'
        f'<span class="ins-txt">{escape(_clip_words(str(h), 20))}</span></div></div>'
        for i, h in enumerate(items)
    )
    return f'<div class="trow ins-row">{cells}</div>'


def _premium_svg_backdrop(P: dict, idx: int, dark: bool = False) -> str:
    ink = "rgba(255,255,255,.42)" if dark else P["primary"]
    accent = P["accent"]
    return f"""
<svg class="premium-backdrop" xmlns="http://www.w3.org/2000/svg" viewBox="0 0 595 842" preserveAspectRatio="none">
  <defs>
    <linearGradient id="premiumGrad{idx}" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0%" stop-color="{P['primary']}" stop-opacity="{'.95' if dark else '.07'}"/>
      <stop offset="62%" stop-color="{P['card2']}" stop-opacity="{'.80' if dark else '.03'}"/>
      <stop offset="100%" stop-color="{accent}" stop-opacity="{'.55' if dark else '.11'}"/>
    </linearGradient>
  </defs>
  <rect width="595" height="842" fill="url(#premiumGrad{idx})"/>
  <path d="M-40 720 C130 590 220 735 380 592 C492 492 548 552 650 414" fill="none" stroke="{accent}" stroke-width="2.5" opacity=".35"/>
  <path d="M-50 112 C120 168 206 52 360 118 C485 172 535 112 650 42" fill="none" stroke="{ink}" stroke-width=".8" opacity=".18"/>
  <circle cx="512" cy="118" r="168" fill="none" stroke="{accent}" stroke-width="1" opacity=".18"/>
  <circle cx="78" cy="752" r="236" fill="{accent}" opacity="{'.10' if dark else '.055'}"/>
</svg>"""


# ─── Main renderer ────────────────────────────────────────────────────────────

def _real_stats(pg: dict) -> list:
    return [s for s in (pg.get("stats") or [])
            if isinstance(s, dict) and str(s.get("number") or s.get("value") or "").strip()]


def _page_has_data(pg: dict) -> bool:
    return bool(
        len(_safe_chart_data(pg.get("chart_data") or [])) >= 2
        or (pg.get("data_table") or {}).get("rows")
        or (pg.get("hero_stat") or {}).get("value")
        or (pg.get("hero_stat") or {}).get("number")
    )


def _vary_pdf_layouts(pages: list, seed: int) -> list:
    """Assign each interior page a *content-faithful* archetype (so nothing the AI
    produced is ever dropped), while injecting per-run variety through mirror flips,
    geometry seeds, and feature↔split alternation on data pages that support both."""
    rng = random.Random(seed)
    n = len(pages)

    def primary(pg: dict) -> str:
        t = str(pg.get("type", "")).lower()
        stats = _real_stats(pg)
        if pg.get("steps"):
            return "timeline"
        if len(stats) >= 2:
            return "statgrid"
        if t in ("case_study",) or pg.get("attribution"):
            return "quote"
        if _page_has_data(pg):
            return "feature"
        if t in ("future", "manifesto") and (pg.get("callout") or pg.get("statements")):
            return "quote"
        if len(pg.get("tiles") or []) >= 3:
            return "cards"
        if len([h for h in (pg.get("highlights") or []) if str(h).strip()]) >= 3 and not pg.get("body"):
            return "cards"
        if pg.get("callout") and not pg.get("body"):
            return "quote"
        return "split"

    prev = None
    for i, pg in enumerate(pages):
        t = str(pg.get("type", "")).lower()
        pg["_geo"] = rng.randrange(4)
        pg["_mirror"] = rng.random() < 0.5
        if i == 0 or t == "cover":
            pg["_layout"] = "cover"; prev = "cover"; continue
        if i == n - 1 or t == "closing":
            pg["_layout"] = "closing"; prev = "closing"; continue
        arche = primary(pg)
        # feature and split both fully present a data+narrative page — alternate for variety.
        if arche in ("feature", "split") and pg.get("body") and _page_has_data(pg):
            if prev in ("feature", "split"):
                arche = "split" if prev == "feature" else "feature"
            elif rng.random() < 0.42:
                arche = "split" if arche == "feature" else "feature"
        pg["_layout"] = arche
        prev = arche
    return pages


def render_pdf_html(structure: dict, page_count: int, font_face_css: str = "") -> str:
    """Premium editorial PDF renderer — WeasyPrint-safe (block/float/table only, no CSS Grid)."""
    P = _resolve_design_tokens(structure)
    title = escape(_clip_chars(structure.get("title") or "Report", 90))
    subtitle = escape(_clip_chars(structure.get("subtitle") or "", 155))
    author = escape(_clip_chars(structure.get("author") or "Prepared by AI Document Service", 90))
    aesthetic = structure.get("aesthetic_direction") or {}
    aesthetic_label = escape(_clip_chars(aesthetic.get("label") or aesthetic.get("name") or "Editorial Intelligence", 42))
    signature = escape(_clip_chars(P.get("signature_element") or "calibrated signal line", 70))
    # Real typography is inlined as base64 @font-face rules (see _build_embedded_font_css).
    # Passed in pre-fetched so the AI-chosen fonts actually render instead of falling back.
    fonts_link = font_face_css or ""

    css = f"""
@page {{ size: A4; margin: 0; }}
* {{ box-sizing: border-box; }}
html, body {{ margin:0; padding:0; font-family:var(--font-body); color:var(--text); }}
:root {{
  --primary:{P['primary']}; --accent:{P['accent']}; --mid:{P['mid']}; --light:{P['light']};
  --wash:{P['wash']}; --text:{P['text']}; --dark:{P['dark']}; --body:{P['body']};
  --card2:{P['card2']}; --muted:{P['muted']};
  --font-display:'{P['h1_font']}', Georgia, 'Times New Roman', serif;
  --font-heading:'{P['h2_font']}', 'Helvetica Neue', Arial, sans-serif;
  --font-body:'{P['body_font']}', 'Helvetica Neue', Arial, sans-serif;
  --font-data:'{P['data_font']}', 'SF Mono', 'Courier New', monospace;
}}
.page {{ width:210mm; height:297mm; page-break-after:always; break-after:page; position:relative; overflow:hidden; background:var(--wash); display:block; }}
.page:last-child {{ page-break-after:auto; break-after:auto; }}
.premium-backdrop {{ position:absolute; inset:0; width:210mm; height:297mm; z-index:0; display:block; }}
.content {{ position:relative; z-index:1; display:block; }}
.kicker {{ font-family:var(--font-data); font-size:7.5pt; letter-spacing:.28em; text-transform:uppercase; font-weight:700; color:var(--accent); display:block; }}
.kicker b {{ color:var(--muted); font-weight:700; }}
h1,h2,h3 {{ margin:0; font-family:var(--font-display); letter-spacing:-.02em; line-height:1.03; display:block; font-weight:800; }}
.cover h1 {{ bookmark-level:1; bookmark-label:content(text); }}
.spread h2, .chapter h2, .dark-page h2 {{ bookmark-level:1; bookmark-label:content(text); }}
p {{ margin:0; color:var(--body); font-size:10pt; line-height:1.62; display:block; }}

/* ── running foot ── */
.folio {{ position:absolute; left:20mm; right:20mm; bottom:11mm; z-index:2; display:block; overflow:hidden; padding-top:3mm; border-top:.35mm solid rgba(0,0,0,.10); font:7pt var(--font-data); color:rgba(0,0,0,.42); letter-spacing:.14em; text-transform:uppercase; }}
.folio.dark {{ border-top-color:rgba(255,255,255,.18); color:rgba(255,255,255,.5); }}
.folio span {{ float:left; display:block; max-width:72%; overflow:hidden; }}
.folio b {{ float:right; display:block; color:var(--accent); font-size:8pt; letter-spacing:.1em; }}

/* ── shared header block ── */
.doc {{ padding:19mm 20mm 20mm; display:block; }}
.doc-head {{ display:block; margin-bottom:8mm; }}
.doc-head .num {{ display:inline-block; font:800 12pt var(--font-data); color:var(--accent); letter-spacing:-.02em; margin-right:4mm; vertical-align:middle; }}
.doc-head h2 {{ font-size:32pt; color:var(--primary); max-width:172mm; margin:4mm 0 0; }}
.lede {{ max-width:158mm; margin-top:6mm; font-size:11pt; line-height:1.62; color:var(--body); }}
.accent-rule {{ display:block; width:22mm; height:1.4mm; background:var(--accent); border-radius:2mm; margin:5mm 0 0; }}

/* ── utility table rows (equal cols, never wrap) ── */
.trow {{ display:table; width:100%; table-layout:fixed; border-collapse:separate; }}
.trow > .cell {{ display:table-cell; vertical-align:top; padding-right:5mm; }}
.trow > .cell:last-child {{ padding-right:0; }}

/* ── cover ── */
.cover {{ background:var(--dark); color:#fff; }}
.cover .content {{ padding:24mm 24mm; height:297mm; position:relative; display:block; }}
.cover-top {{ margin-bottom:12mm; display:block; }}
.cover-pill {{ display:inline-block; border:1px solid rgba(255,255,255,.24); border-radius:999px; padding:2.4mm 7mm; background:rgba(255,255,255,.07); margin-right:4mm; font:7pt var(--font-data); letter-spacing:.24em; text-transform:uppercase; color:rgba(255,255,255,.78); vertical-align:middle; }}
.cover-body {{ display:block; min-height:120mm; }}
.cover h1 {{ max-width:168mm; font-size:56pt; color:#fff; line-height:1.0; }}
.cover-sub {{ max-width:140mm; margin-top:10mm; color:rgba(255,255,255,.8); font-size:14pt; line-height:1.5; font-family:var(--font-heading); }}
.cover-callout {{ max-width:130mm; color:#fff; font:italic 15pt var(--font-display); line-height:1.4; }}
.cover-meta {{ position:absolute; left:24mm; right:24mm; bottom:22mm; display:block; overflow:hidden; }}
.cover-meta > div:first-child {{ float:left; max-width:calc(100% - 56mm); display:block; }}
.cover-num {{ float:right; display:block; font:800 78pt var(--font-data); color:rgba(255,255,255,.12); letter-spacing:-.08em; line-height:.8; }}
.cover-author {{ margin-top:8mm; color:rgba(255,255,255,.5); font:7pt var(--font-data); letter-spacing:.2em; text-transform:uppercase; }}

/* ── narrative / body ── */
.body-col p + p {{ margin-top:4mm; }}
.pullquote {{ margin-top:8mm; padding:2mm 0 2mm 8mm; border-left:1.4mm solid var(--accent); font:italic 15pt var(--font-display); line-height:1.4; color:var(--primary); display:block; }}
.dark-page .pullquote {{ color:#fff; border-color:var(--accent); }}
.insights {{ margin-top:7mm; display:block; }}
.insight {{ display:block; position:relative; padding:3.6mm 0 3.6mm 6mm; border-top:.3mm solid rgba(0,0,0,.10); font-size:9.4pt; line-height:1.42; color:var(--body); }}
.insight:before {{ content:""; position:absolute; left:0; top:4.6mm; width:2.6mm; height:2.6mm; border-radius:50%; background:var(--accent); }}
.dark-page .insight {{ color:rgba(255,255,255,.78); border-color:rgba(255,255,255,.14); }}

/* ── insight cards row ── */
.ins-row {{ margin-top:8mm; }}
.ins-card {{ display:block; background:#fff; border:.3mm solid rgba(15,23,42,.09); border-top:1.4mm solid var(--accent); border-radius:2.4mm; padding:5.5mm 5mm; min-height:33mm; box-shadow:0 10px 26px rgba(15,23,42,.07); }}
.ins-num {{ display:block; font:800 11pt var(--font-data); color:var(--accent); letter-spacing:-.02em; margin-bottom:3.5mm; }}
.ins-txt {{ display:block; color:#334155; font:9.6pt/1.5 var(--font-body); }}
.dark-page .ins-card {{ background:rgba(255,255,255,.05); border-color:rgba(255,255,255,.12); box-shadow:none; }}
.dark-page .ins-txt {{ color:rgba(255,255,255,.82); }}
/* takeaway/cards archetype: taller cards that fill the page */
.cards-hero {{ margin-top:11mm; }}
.cards-hero .ins-card {{ min-height:132mm; padding:9mm 7mm; border-top-width:2mm; }}
.cards-hero .ins-num {{ font-size:15pt; margin-bottom:8mm; }}
.cards-hero .ins-txt {{ font-size:11pt; line-height:1.55; }}

/* ── two-column split (table layout → guaranteed top-alignment) ── */
.split {{ margin-top:2mm; display:table; width:100%; table-layout:fixed; }}
.split .main-col {{ display:table-cell; vertical-align:top; }}
.split .side-col {{ display:table-cell; vertical-align:top; width:78mm; }}
.split .gut {{ padding-right:9mm; }}

/* ── hero stat card ── */
.premium-stat {{ display:block; background:linear-gradient(140deg,var(--primary),var(--card2) 60%,var(--accent)); color:#fff; padding:11mm; border-radius:3.5mm; box-shadow:0 16px 40px rgba(15,23,42,.24); min-height:98mm; position:relative; overflow:hidden; }}
.premium-stat:after {{ content:""; position:absolute; right:-30mm; top:-30mm; width:100mm; height:100mm; border-radius:50%; border:1.2mm solid rgba(255,255,255,.12); }}
.premium-stat-value {{ display:block; margin-top:20mm; font:800 72pt var(--font-data); letter-spacing:-.05em; line-height:.9; }}
.premium-stat-label {{ display:block; margin-top:6mm; font:800 9pt var(--font-data); letter-spacing:.2em; text-transform:uppercase; color:rgba(255,255,255,.82); }}
.premium-stat p {{ margin-top:5mm; max-width:120mm; color:rgba(255,255,255,.76); font-size:10pt; line-height:1.55; }}
.side-col .premium-stat {{ min-height:150mm; padding:10mm; }}
.side-col .premium-stat-value {{ margin-top:58mm; font-size:54pt; }}

/* ── chart ── */
.premium-chart {{ display:block; background:linear-gradient(160deg,#14161c,#0c0d11); color:#fff; padding:9mm 11mm; border-radius:3.5mm; box-shadow:0 16px 40px rgba(15,23,42,.22); }}
.premium-chart.big {{ min-height:90mm; }}
.side-col .premium-chart {{ min-height:150mm; }}
.chart-head {{ display:block; overflow:hidden; margin-bottom:7mm; padding-bottom:4mm; border-bottom:.3mm solid rgba(255,255,255,.12); }}
.chart-kicker {{ float:left; font:800 8.5pt var(--font-data); letter-spacing:.16em; text-transform:uppercase; color:rgba(255,255,255,.82); max-width:70%; }}
.chart-delta {{ float:right; font:800 8.5pt var(--font-data); color:var(--accent); letter-spacing:.06em; }}
.chart-rows {{ display:block; }}
.chart-row {{ display:table; width:100%; table-layout:fixed; margin:5.2mm 0; border-collapse:separate; border-spacing:0; }}
.chart-label {{ display:table-cell; width:30mm; vertical-align:middle; padding-right:3mm; font:8pt var(--font-data); color:rgba(255,255,255,.72); }}
.chart-track {{ display:table-cell; vertical-align:middle; height:6mm; background:rgba(255,255,255,.08); border-radius:999px; overflow:hidden; }}
.chart-value {{ display:table-cell; width:16mm; vertical-align:middle; text-align:right; font:9pt var(--font-data); color:#fff; font-weight:800; }}
.chart-bar {{ display:block; height:6mm; background:linear-gradient(90deg,var(--mid),var(--accent)); border-radius:999px; }}
.chart-bar.peak {{ background:linear-gradient(90deg,var(--accent),#fff); }}

/* ── table ── */
.hero-table {{ display:block; border-radius:3.5mm; overflow:hidden; box-shadow:0 14px 34px rgba(15,23,42,.12); }}
.premium-table {{ display:table; width:100%; border-collapse:collapse; overflow:hidden; font-size:9pt; background:#fff; }}
.premium-table thead {{ display:table-header-group; }}
.premium-table tbody {{ display:table-row-group; }}
.premium-table tr {{ display:table-row; }}
.premium-table th {{ display:table-cell; background:var(--primary); color:#fff; text-align:left; padding:4.5mm 4mm; font:800 7.5pt var(--font-data); letter-spacing:.1em; text-transform:uppercase; }}
.premium-table td {{ display:table-cell; padding:4.2mm 4mm; border-bottom:.25mm solid rgba(15,23,42,.08); color:#243042; }}
.premium-table tbody tr:first-child td {{ font-weight:700; }}
.premium-table tr:nth-child(even) td {{ background:var(--wash); }}

/* ── photo ── */
.photo-plate {{ display:block; margin:0; border-radius:3.5mm; overflow:hidden; position:relative; box-shadow:0 16px 40px rgba(15,23,42,.22); background:#111; }}
.photo-plate img {{ width:100%; height:96mm; object-fit:cover; display:block; }}
.photo-plate.big img {{ height:96mm; }}
.side-col .photo-plate img {{ height:150mm; }}
.photo-plate:after {{ content:""; position:absolute; inset:0; background:linear-gradient(180deg,transparent 50%,rgba(0,0,0,.5)); }}
.photo-caption {{ position:absolute; left:6mm; right:6mm; bottom:5mm; z-index:2; color:#fff; font:7.5pt var(--font-data); letter-spacing:.14em; text-transform:uppercase; display:block; }}

/* ── geometric data panel ── */
.geo-panel {{ display:block; border-radius:3.5mm; overflow:hidden; box-shadow:0 16px 40px rgba(15,23,42,.2); line-height:0; }}
.geo-panel svg {{ display:block; width:100%; height:96mm; }}
.side-col .geo-panel svg {{ height:150mm; }}

/* ── big stat grid ── */
.stat-row {{ margin-top:9mm; }}
.bigstat {{ display:block; background:#fff; border:.3mm solid rgba(15,23,42,.08); border-radius:3.5mm; padding:9mm 7mm; min-height:150mm; position:relative; overflow:hidden; box-shadow:0 12px 30px rgba(15,23,42,.08); }}
.bigstat:before {{ content:""; position:absolute; left:0; top:0; width:100%; height:2mm; background:var(--accent); }}
.bigstat-idx {{ display:block; font:800 8pt var(--font-data); color:var(--muted); letter-spacing:.16em; }}
.bigstat-val {{ display:block; margin-top:46mm; font:800 46pt var(--font-display); color:var(--primary); letter-spacing:-.03em; line-height:.92; }}
.bigstat-lbl {{ display:block; margin-top:6mm; font:800 8pt var(--font-data); letter-spacing:.16em; text-transform:uppercase; color:var(--accent); }}
.bigstat-ctx {{ display:block; margin-top:5mm; font:9.6pt/1.55 var(--font-body); color:#475569; }}

/* ── timeline ── */
.tl-row {{ margin-top:9mm; }}
.tl-card {{ display:block; background:#fff; border:.3mm solid rgba(15,23,42,.08); border-radius:3.5mm; padding:8mm 6mm; min-height:148mm; position:relative; box-shadow:0 12px 30px rgba(15,23,42,.07); }}
.tl-node {{ display:block; width:9mm; height:9mm; border-radius:50%; background:var(--accent); color:#fff; text-align:center; font:800 8pt/9mm var(--font-data); margin-bottom:6mm; }}
.tl-date {{ display:block; font:800 12pt var(--font-data); color:var(--primary); letter-spacing:-.02em; }}
.tl-step {{ display:block; margin-top:2mm; font:800 8pt var(--font-data); letter-spacing:.14em; text-transform:uppercase; color:var(--accent); }}
.tl-desc {{ display:block; margin-top:4mm; font:9.6pt/1.55 var(--font-body); color:#334155; }}

/* ── dark feature (quote / manifesto) ── */
.dark-page {{ background:var(--dark); color:#fff; }}
.dark-page .content {{ padding:22mm 22mm 20mm; display:block; }}
.dark-page h2 {{ color:#fff; }}
.dark-page p {{ color:rgba(255,255,255,.76); }}
.dark-page .kicker {{ color:var(--muted); }}
.dark-page .doc-head h2 {{ color:#fff; }}
.giant-quote {{ display:block; margin:14mm 0 0; font:italic 34pt/1.22 var(--font-display); color:#fff; max-width:170mm; letter-spacing:-.01em; }}
.giant-quote .mark {{ font-size:52pt; color:var(--accent); line-height:0; }}
.quote-attr {{ display:block; margin-top:10mm; font:800 10pt var(--font-data); letter-spacing:.12em; text-transform:uppercase; color:var(--accent); }}

/* ── closing ── */
.closing h2 {{ font-size:44pt; }}
"""

    pages = []
    source_pages = structure.get("pages", [])[:page_count]
    # Per-document randomised, content-aware layout assignment (different every run).
    _vary_pdf_layouts(source_pages, random.randrange(1 << 30))
    for idx, page in enumerate(source_pages, 1):
        ptype = page.get("_layout") or "split"
        mir = " mirror" if page.get("_mirror") else ""
        eyebrow = escape(_clip_chars(page.get("eyebrow") or f"Section {idx}", 44))
        heading = _clip_words(page.get("heading") or structure.get("title") or "Key Insight", 13)
        heading_html = escape(heading)
        body = escape(_clip_words(page.get("body") or "", 96 if idx in (1, page_count) else 82))
        callout = escape(_clip_words(page.get("callout") or "", 30))
        highlights_raw = [str(h) for h in (page.get("highlights") or []) if str(h).strip()][:4]
        highlights = [escape(_clip_words(h, 18)) for h in highlights_raw]
        dark = ptype in ("cover", "quote", "manifesto")
        folio_cls = "folio dark" if dark else "folio"
        folio = f'<div class="{folio_cls}"><span>{title}</span><b>{idx:02d} / {page_count:02d}</b></div>'
        head_html = (f'<div class="doc-head"><div class="kicker"><span class="num">{idx:02d}</span>{eyebrow}</div>'
                     f'<h2>{heading_html}</h2><span class="accent-rule"></span></div>')
        pullquote = f'<div class="pullquote">{callout}</div>' if callout else ""
        insights_list = "".join(f'<div class="insight">{h}</div>' for h in highlights)

        if ptype == "cover":
            pages.append(f"""
<section class="page cover">
  {_premium_svg_backdrop(P, idx, dark=True)}
  <div class="content">
    <div class="cover-top"><span class="cover-pill">{aesthetic_label}</span><span class="cover-pill">{eyebrow}</span></div>
    <div class="cover-body"><h1>{heading_html}</h1><p class="cover-sub">{subtitle}</p></div>
    <div class="cover-meta"><div><p class="cover-callout">{callout or subtitle}</p><p class="cover-author">{author} &nbsp;·&nbsp; motif: {signature}</p></div><div class="cover-num">01</div></div>
  </div>
</section>""")

        elif ptype in ("quote", "manifesto"):
            statements = [str(s) for s in (page.get("statements") or []) if str(s).strip()]
            chips = _insight_cards(statements or highlights_raw, P, limit=3)
            attr = escape(_clip_chars(page.get("attribution") or "", 48))
            quote_text = callout or body
            pages.append(f"""
<section class="page dark-page">
  {_premium_svg_backdrop(P, idx, dark=True)}
  <div class="content">
    <div class="doc-head"><div class="kicker"><span class="num" style="color:var(--accent)">{idx:02d}</span>{eyebrow}</div><h2>{heading_html}</h2></div>
    <div class="giant-quote"><span class="mark">“</span>{quote_text}</div>
    {f'<div class="quote-attr">{attr}</div>' if attr else ''}
    {chips}
  </div>
  {folio}
</section>""")

        elif ptype == "closing":
            takeaways = [str(t) for t in (page.get("takeaways") or page.get("statements") or page.get("highlights") or [])]
            cards = _insight_cards(takeaways, P, limit=3)
            pages.append(f"""
<section class="page closing" style="background:var(--wash);">
  {_premium_svg_backdrop(P, idx)}
  <div class="content doc">
    <div class="doc-head"><div class="kicker"><span class="num">{idx:02d}</span>{eyebrow}</div><h2>{heading_html}</h2><span class="accent-rule"></span></div>
    <p class="lede">{body}</p>
    {pullquote}
    {cards}
  </div>
  {folio}
</section>""")

        elif ptype == "statgrid":
            stats = _real_stats(page)[:3]
            cells = "".join(f'<div class="cell">{_big_stat_card(s, P, i)}</div>' for i, s in enumerate(stats))
            pages.append(f"""
<section class="page">
  {_premium_svg_backdrop(P, idx)}
  <div class="content doc">
    {head_html}
    <p class="lede">{body}</p>
    <div class="trow stat-row">{cells}</div>
  </div>
  {folio}
</section>""")

        elif ptype == "timeline":
            steps = page.get("steps") or []
            if not steps:
                steps = [{"step": h, "desc": "", "date": f"{i+1:02d}"} for i, h in enumerate(highlights_raw)]
            steps = steps[:4]
            cells = ""
            for i, s in enumerate(steps):
                date = escape(_clip_chars(s.get("date") or f"{i+1:02d}", 16))
                step = escape(_clip_chars(s.get("step") or f"Phase {i+1}", 22))
                desc = escape(_clip_words(s.get("desc") or s.get("step") or "", 24))
                cells += (f'<div class="cell"><div class="tl-card"><span class="tl-node">{i+1}</span>'
                          f'<span class="tl-date">{date}</span><span class="tl-step">{step}</span>'
                          f'<span class="tl-desc">{desc}</span></div></div>')
            pages.append(f"""
<section class="page">
  {_premium_svg_backdrop(P, idx)}
  <div class="content doc">
    {head_html}
    <p class="lede">{body}</p>
    <div class="trow tl-row">{cells}</div>
  </div>
  {folio}
</section>""")

        elif ptype == "cards":
            tiles = page.get("tiles") or []
            if tiles:
                items = [f"{t.get('value','')}".strip() or f"{t.get('label','')}" for t in tiles[:4]]
                labels = [str(t.get("label", "") or "").strip() for t in tiles[:4]]
            else:
                items = highlights_raw[:4]
                labels = ["" for _ in items]
            cells = "".join(
                f'<div class="cell"><div class="ins-card"><span class="ins-num">{escape(str(labels[i] or f"{i+1:02d}"))}</span>'
                f'<span class="ins-txt">{escape(_clip_words(str(it), 26))}</span></div></div>'
                for i, it in enumerate(items)
            )
            pages.append(f"""
<section class="page">
  {_premium_svg_backdrop(P, idx)}
  <div class="content doc">
    {head_html}
    <p class="lede">{body}</p>
    <div class="trow ins-row cards-hero">{cells}</div>
  </div>
  {folio}
</section>""")

        elif ptype == "feature":
            hero = _hero_visual_html(page, P, idx, heading)
            cards = _insight_cards(highlights_raw, P, limit=3)
            pages.append(f"""
<section class="page">
  {_premium_svg_backdrop(P, idx)}
  <div class="content doc">
    {head_html}
    <p class="lede">{body}</p>
    <div style="margin-top:8mm">{hero}</div>
    {cards}
  </div>
  {folio}
</section>""")

        else:  # split — narrative + tall side visual
            visual = _side_visual_html(page, P, idx, heading)
            narrative = f'<p>{body}</p>{pullquote}<div class="insights">{insights_list}</div>'
            if page.get("_mirror"):  # visual on the left, narrative on the right
                cols = (f'<div class="side-col gut">{visual}</div>'
                        f'<div class="main-col body-col">{narrative}</div>')
            else:  # narrative on the left, visual on the right
                cols = (f'<div class="main-col body-col gut">{narrative}</div>'
                        f'<div class="side-col">{visual}</div>')
            pages.append(f"""
<section class="page">
  {_premium_svg_backdrop(P, idx)}
  <div class="content doc">
    {head_html}
    <div class="split">{cols}</div>
  </div>
  {folio}
</section>""")

    pagination_css = CSS(string="""
        @page { size: A4; margin: 0; }
        .page { page-break-after: always; break-after: page; }
        .page:last-child { page-break-after: auto; break-after: auto; }
        .premium-stat, .premium-chart, .premium-table, .photo-plate, .geo-panel, .bigstat, .tl-card, .ins-card, .pullquote { break-inside: avoid; page-break-inside: avoid; }
        h1, h2, h3 { break-after: avoid; page-break-after: avoid; }
        p { orphans: 3; widows: 3; }
    """)
    html_doc = f"""<!DOCTYPE html><html lang="en"><head><meta charset="utf-8"/>{fonts_link}<style>{css}</style></head><body>{"".join(pages)}</body></html>"""
    return html_doc, pagination_css


def _quality_check(pdf_bytes: bytes) -> dict:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(pdf_bytes))
        pages = len(reader.pages)
        size_kb = len(pdf_bytes) / 1024
        issues = []
        if pages < 8:
            issues.append(f"Warning: Only {pages} pages generated (fewer than 8).")
        if size_kb < 80:
            issues.append(f"Warning: File size is {size_kb:.1f}KB (under 80KB, styles may be missing).")
        elif size_kb > 15000:
            issues.append(f"Warning: File size is {size_kb:.1f}KB (over 15MB, unoptimised).")
        
        # Check dimensions of first page (A4 is approx 595x842 points)
        if pages > 0:
            box = reader.pages[0].mediabox
            w, h = float(box.width), float(box.height)
            if not (590 <= w <= 600 and 837 <= h <= 847) and not (837 <= w <= 847 and 590 <= h <= 600):
                issues.append(f"Warning: Invalid page dimensions {w}x{h} (not A4).")
                
        return {"pages": pages, "size_kb": round(size_kb, 1), "issues": issues, "passed": len(issues) == 0}
    except Exception as e:
        return {"pages": 0, "size_kb": 0, "issues": [f"Quality check failed: {str(e)}"], "passed": False}


@app.post("/docs/generate/pdf")
async def generate_pdf(payload: dict):
    try:
        prompt = _normalize_topic_prompt(payload["prompt"])
        page_count = requested_count(prompt, "page", default=10, minimum=8, maximum=12)
        structure = await call_ai_pdf(prompt, page_count)
        # Fetch AI-chosen fonts and topic images concurrently (both time-boxed).
        tokens = _resolve_design_tokens(structure)
        font_css, _ = await asyncio.gather(
            _build_embedded_font_css(tokens),
            _prefetch_images(structure),
        )
        html_doc, pagination_css = render_pdf_html(structure, page_count, font_face_css=font_css)
        buffer = io.BytesIO()
        HTML(string=html_doc).write_pdf(buffer, stylesheets=[pagination_css])
        pdf_bytes = buffer.getvalue()
        
        # Run quality check
        qa_result = _quality_check(pdf_bytes)
        print(f"PDF QA Check: {qa_result}")
        
        buffer.seek(0)
        blob = await upload_to_vercel_blob(buffer, "output.pdf", "application/pdf")
        blob["qa_check"] = qa_result
        return blob
    except Exception as exc:
        _generation_error(exc)


# ══════════════════════════════════════════════════════════════════════════════
# PPTX GENERATION
# ══════════════════════════════════════════════════════════════════════════════

async def call_ai_pptx(prompt: str, slide_count: int) -> dict:
    core = _creative_director_core("slide", slide_count, "slide")
    system = f"""{core}

Generate exactly {slide_count} presentation slides. slides[0] is ALWAYS the title/cover slide.

SLIDE TYPES (narrative order — each uses a distinct layout):
- "title"       → slide 1: cinematic cover, aesthetic declared in subtitle/tagline
- "context"     → overview with 3 stat cards establishing stakes
- "deep_dive_1" → two-column comparison of key factors
- "deep_dive_2" → image_text: body + highlight sidebar, visual storytelling
- "deep_dive_3" → 3 stat cards with real quantitative data
- "comparative" → two_col: before/after or side-by-side analysis
- "case_study"  → pull quote + attribution + 3 supporting highlight cards
- "future"      → image_text: forward-looking body + aspirational highlights
- "takeaways"   → 3 stat-style insight cards — NOT bullet list
- "closing"     → final slide: headline + CTA + signature tagline

JSON schema:
{{
  "title": "Deck title",
  "subtitle": "Subtitle line",
  "aesthetic_direction": {{
    "name": "bauhaus_data",
    "label": "Bauhaus Data Art",
    "rationale": "Why this aesthetic fits (2 sentences)"
  }},
  "design_system": {{
    "colors": {{
      "dominant": "#1A1A2E", "secondary": "#16213E",
      "supporting_1": "#0F3460", "supporting_2": "#E8E8E8", "accent": "#E94560"
    }},
    "typography": {{
      "h1_font": "Archivo Black", "h2_font": "Work Sans",
      "body_font": "IBM Plex Sans", "data_font": "Space Mono"
    }},
    "signature_element": "Geometric triangle accent in top-right corner"
  }},
  "palette": "midnight",
  "slides": [
    {{
      "type": "title|context|deep_dive_1|deep_dive_2|deep_dive_3|comparative|case_study|future|takeaways|closing",
      "title": "Slide headline",
      "subtitle": "Title slide only",
      "tagline": "Title or closing — include aesthetic direction label",
      "bullets": [],
      "left_points": ["Left column point with 15+ words of detail"],
      "right_points": ["Right column point with 15+ words of detail"],
      "stats": [{{"number":"847M","label":"METRIC","context":"Real context sentence"}}],
      "quote": "Full memorable quote",
      "attribution": "— Source Name, Year",
      "body": "Rich paragraph for image_text slides, 60+ words",
      "highlights": ["Sidebar insight with real data"],
      "chart_data": [["Label", 42]],
      "cta": "Call to action for closing slide"
    }}
  ]
}}

RULES:
- slides[0].type MUST be "title"
- slides[{slide_count - 1}].type MUST be "closing"
- Follow narrative order; skip middle sections only if slide_count < 10
- stat/context/deep_dive_3/takeaways slides: exactly 3 stat objects with REAL numbers
- two_col/comparative/deep_dive_1: 3 items in left_points AND 3 in right_points
- case_study: quote + attribution + 3 highlights
- image_text/deep_dive_2/future: body 60+ words + 3-4 highlights
- NO bullet-list-primary slides — use stats, columns, or quote layouts
- slides array MUST contain exactly {slide_count} items"""

    raw = await _deepinfra_call([
        {"role": "system", "content": system},
        {"role": "user", "content": prompt},
    ], temperature=0.4)
    data = _as_json(raw, "Presentation content generator")
    slides = data.get("slides", [])
    while len(slides) < slide_count:
        i = len(slides)
        seq_idx = min(i, len(PPTX_SLIDE_SEQUENCE) - 1)
        lt = PPTX_SLIDE_SEQUENCE[seq_idx]
        slides.append({
            "type": lt, "title": f"Key Insight {i+1}",
            "stats": [{"number": "—", "label": "Metric", "context": "Supporting context for this insight."}] * 3,
            "left_points": ["Left perspective with actionable detail on the topic."] * 3,
            "right_points": ["Right perspective with complementary analysis."] * 3,
            "highlights": ["Key insight with real relevance to the topic."] * 3,
            "body": "This slide presents essential detail on the topic with clear, actionable perspective crafted for maximum audience value.",
        })
    data["slides"] = slides[:slide_count]
    if slide_count >= 2:
        data["slides"][0]["type"] = "title"
        data["slides"][-1]["type"] = "closing"
    return data


def _shade(c: tuple, f: float) -> tuple:
    """Lighten (f>1) or darken (f<1) an RGB tuple, clamped to 0-255."""
    return tuple(max(0, min(255, int(round(v * f)))) for v in c)


def _hex_to_rgb_tuple(hex_color: str) -> tuple:
    h = str(hex_color).lstrip("#")
    if len(h) != 6:
        return (0x1E, 0x27, 0x61)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _pptx_palette_from_structure(structure: dict) -> dict:
    tokens = _resolve_design_tokens(structure)
    bg = _hex_to_rgb_tuple(tokens["primary"])
    card = _hex_to_rgb_tuple(tokens["card2"])
    accent = _hex_to_rgb_tuple(tokens["accent"])
    light = _hex_to_rgb_tuple(tokens["light"])
    white = (0xFF, 0xFF, 0xFF)
    gray = _hex_to_rgb_tuple(tokens["mid"])
    return {
        "bg": bg, "primary": bg, "accent": accent, "light": light,
        "white": white, "gray": gray, "card": card,
        "h1_font": tokens["h1_font"], "h2_font": tokens["h2_font"],
        "body_font": tokens["body_font"], "data_font": tokens["data_font"],
    }


PPTX_PALETTES = {
    "midnight":  {"bg": (0x1E,0x27,0x61), "primary": (0x1E,0x27,0x61), "accent": (0x49,0x8F,0xD6), "light": (0xCA,0xDC,0xFC), "white": (0xFF,0xFF,0xFF), "gray": (0x8A,0x9A,0xBF), "card": (0x28,0x34,0x76)},
    "forest":    {"bg": (0x2C,0x5F,0x2D), "primary": (0x2C,0x5F,0x2D), "accent": (0x97,0xBC,0x62), "light": (0xD4,0xED,0xC1), "white": (0xFF,0xFF,0xFF), "gray": (0x7A,0x9A,0x5A), "card": (0x36,0x72,0x37)},
    "coral":     {"bg": (0xF9,0x61,0x67), "primary": (0x2F,0x3C,0x7E), "accent": (0xF9,0x61,0x67), "light": (0xFD,0xE0,0xE0), "white": (0xFF,0xFF,0xFF), "gray": (0x88,0x77,0xAA), "card": (0x3A,0x4A,0x9A)},
    "terracotta":{"bg": (0xB8,0x50,0x42), "primary": (0xB8,0x50,0x42), "accent": (0xE7,0xE8,0xD1), "light": (0xF5,0xF0,0xE8), "white": (0xFF,0xFF,0xFF), "gray": (0xA7,0xBE,0xAE), "card": (0xD4,0x62,0x52)},
    "ocean":     {"bg": (0x06,0x5A,0x82), "primary": (0x06,0x5A,0x82), "accent": (0x02,0xC3,0x9A), "light": (0xC8,0xF0,0xE8), "white": (0xFF,0xFF,0xFF), "gray": (0x7A,0xB8,0xD4), "card": (0x1C,0x72,0x93)},
    "charcoal":  {"bg": (0x36,0x45,0x4F), "primary": (0x36,0x45,0x4F), "accent": (0x7A,0xC0,0xD4), "light": (0xE8,0xF2,0xF5), "white": (0xFF,0xFF,0xFF), "gray": (0x8A,0x9A,0xA8), "card": (0x48,0x58,0x64)},
    "berry":     {"bg": (0x6D,0x2E,0x46), "primary": (0x6D,0x2E,0x46), "accent": (0xEC,0xE2,0xD0), "light": (0xF5,0xEC,0xE5), "white": (0xFF,0xFF,0xFF), "gray": (0xA2,0x67,0x69), "card": (0x84,0x3A,0x58)},
    "cherry":    {"bg": (0x99,0x00,0x11), "primary": (0x99,0x00,0x11), "accent": (0xFC,0xF6,0xF5), "light": (0xFF,0xE8,0xE8), "white": (0xFF,0xFF,0xFF), "gray": (0xCC,0x88,0x88), "card": (0xB0,0x10,0x20)},
}


def _rgb(t: tuple) -> PPTXRGBColor:
    return PPTXRGBColor(*t)

def _solid(shape, color: PPTXRGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _add_rect(slide, l, t, w, h, color, line=False):
    from pptx.util import Inches as I, Pt
    shape = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    _solid(shape, color)
    if not line:
        shape.line.fill.background()
    return shape


def _grad(shape, c1: tuple, c2: tuple, angle: float = 90.0):
    """Apply a two-stop linear gradient fill; falls back to solid on any failure."""
    try:
        f = shape.fill
        f.gradient()
        stops = f.gradient_stops
        stops[0].color.rgb = _rgb(c1)
        stops[1].color.rgb = _rgb(c2)
        try:
            f.gradient_angle = angle
        except Exception:
            pass
    except Exception:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(c1)
    shape.line.fill.background()
    return shape


def _add_grad_rect(slide, l, t, w, h, c1, c2, angle=90.0):
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    return _grad(shape, c1, c2, angle)


def _soft_shadow(shape, blur_pt=13, dist_pt=7, dir_deg=90, alpha_pct=30, color="0B1220"):
    """Inject a soft outer drop shadow (python-pptx has no high-level API for it)."""
    try:
        spPr = shape._element.spPr
        old = spPr.find(pptx_qn('a:effectLst'))
        if old is not None:
            spPr.remove(old)
        blur = int(blur_pt * 12700); dist = int(dist_pt * 12700)
        direction = int(dir_deg * 60000); alpha = int(alpha_pct * 1000)
        xml = (f'<a:effectLst {pptx_nsdecls("a")}>'
               f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{direction}" rotWithShape="0">'
               f'<a:srgbClr val="{color}"><a:alpha val="{alpha}"/></a:srgbClr>'
               f'</a:outerShdw></a:effectLst>')
        spPr.append(pptx_parse_xml(xml))
    except Exception:
        pass
    return shape


def _round_rect(slide, l, t, w, h, color=None, grad=None, shadow=True, radius=0.08):
    """A rounded rectangle card with optional gradient fill + soft shadow."""
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(5, I(l), I(t), I(w), I(h))  # 5 = ROUNDED_RECTANGLE
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    if grad is not None:
        _grad(shape, grad[0], grad[1], grad[2] if len(grad) > 2 else 90.0)
    else:
        _solid(shape, _rgb(color))
        shape.line.fill.background()
    if shadow:
        _soft_shadow(shape)
    return shape

def _txt(slide, text, l, t, w, h, size=16, bold=False, color=None, align=PP_ALIGN.LEFT,
         italic=False, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(PPTXInches(l), PPTXInches(t), PPTXInches(w), PPTXInches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPTXPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return tb

def _add_bullets_textbox(slide, bullets, l, t, w, h, size=14, color=None, marker="•",
                          indent_color=None, font="Calibri"):
    tb = slide.shapes.add_textbox(PPTXInches(l), PPTXInches(t), PPTXInches(w), PPTXInches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = PPTXPt(6)
        run = p.add_run()
        run.text = f"{marker}  {b}"
        run.font.size = PPTXPt(size)
        run.font.name = font
        if color:
            run.font.color.rgb = color


def _pptx_decor(slide, pal, stype):
    """Add subtle decorative corner accents to content slides."""
    ACC = _rgb(pal["accent"])
    CARD = _rgb(pal["card"])
    if stype in ("bullets", "two_col", "image_text"):
        corner = slide.shapes.add_shape(1, PPTXInches(12.4), PPTXInches(0), PPTXInches(0.93), PPTXInches(1.35))
        _solid(corner, ACC)
    elif stype in ("stat", "quote"):
        circ = slide.shapes.add_shape(9, PPTXInches(-0.8), PPTXInches(5.5), PPTXInches(2.8), PPTXInches(2.8))
        circ.fill.solid(); circ.fill.fore_color.rgb = CARD; circ.line.fill.background()


# ─── PPTX motion: entrance animations, slide transitions, interactivity ───────
#
# python-pptx has no animation API, so we inject the OOXML <p:timing> and
# <p:transition> trees directly. Effects use nodeType="afterEffect" so the whole
# slide auto-plays on display (no clicking required), staggered for a build-in feel.

# entrance effect catalog → (animEffect filter, PowerPoint presetID)
_PPTX_ENTRANCES = {
    "fade":      ("fade", 10),
    "wipe_up":   ("wipe(up)", 22),
    "wipe_right":("wipe(right)", 22),
    "wipe_left": ("wipe(left)", 22),
    "blinds":    ("blinds(horizontal)", 3),
    "dissolve":  ("dissolve", 9),
    "circle":    ("circle", 5),
    "diamond":   ("diamond", 7),
    "wheel":     ("wheel(1)", 21),
    "plus":      ("plus", 16),
    "randombar": ("randombar(vertical)", 18),
    "strips":    ("strips(downRight)", 19),
}

# slide-to-slide transitions (one chosen per deck for cohesion)
_PPTX_TRANSITIONS = [
    '<p:fade/>',
    '<p:fade thruBlk="1"/>',
    '<p:push dir="l"/>',
    '<p:push dir="u"/>',
    '<p:cover dir="d"/>',
    '<p:wipe dir="d"/>',
    '<p:split orient="horz" dir="out"/>',
    '<p:zoom dir="in"/>',
    '<p:cut/>',
]


def _anim_effect_par(cid: int, spid: int, delay: int, dur: int, filt: str, preset_id: int) -> str:
    """One staggered entrance effect for a single shape (3-level cTn nest, PowerPoint style)."""
    return (
        f'<p:par><p:cTn id="{cid}" fill="hold">'
        f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
        f'<p:childTnLst><p:par><p:cTn id="{cid+1}" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst><p:par>'
        f'<p:cTn id="{cid+2}" presetID="{preset_id}" presetClass="entr" presetSubtype="0" '
        f'fill="hold" nodeType="afterEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'<p:set><p:cBhvr>'
        f'<p:cTn id="{cid+3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="{filt}"><p:cBhvr>'
        f'<p:cTn id="{cid+4}" dur="{dur}"/>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    )


def _build_pptx_timing(specs: list):
    """specs: list of (spid, delay_ms, dur_ms, filter, preset_id). Returns a <p:timing> element."""
    cid = 3
    pars = []
    for spid, delay, dur, filt, preset in specs:
        pars.append(_anim_effect_par(cid, spid, delay, dur, filt, preset))
        cid += 5
    xml = (
        f'<p:timing {pptx_nsdecls("p", "a")}><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(pars)}</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    return pptx_parse_xml(xml)


def _set_slide_transition(slide, transition_xml: str, advance_after: float = None):
    """Inject a <p:transition>. Must be called BEFORE _animate_slide to keep schema order."""
    adv = ""
    if advance_after is not None:
        adv = f'<p:advTm/>'  # placeholder kept off; manual advance preferred
    xml = f'<p:transition {pptx_nsdecls("p")} spd="med">{transition_xml}</p:transition>'
    slide._element.append(pptx_parse_xml(xml))


def _animate_slide(slide, filt_key: str, max_shapes: int = 9, step: int = 170, dur: int = 460):
    """Auto-playing staggered entrance for a slide's shapes (background shape stays static)."""
    shapes = list(slide.shapes)
    if len(shapes) < 2:
        return
    filt, preset = _PPTX_ENTRANCES.get(filt_key, _PPTX_ENTRANCES["fade"])
    specs = []
    for i, sh in enumerate(shapes[1:max_shapes + 1]):
        specs.append((sh.shape_id, 100 + i * step, dur, filt, preset))
    if specs:
        slide._element.append(_build_pptx_timing(specs))


def _add_nav_controls(slide, prs, idx: int, total: int, pal: dict):
    """Add subtle clickable prev / home / next chevrons that jump between slides."""
    ACC = _rgb(pal["accent"])
    targets = []
    if idx > 0:
        targets.append(("‹", idx - 1))           # prev
    targets.append(("●", 0))                       # home (slide 1)
    if idx < total - 1:
        targets.append(("›", idx + 1))           # next
    bx = 13.33 - 0.05 - len(targets) * 0.42
    for label, tgt in targets:
        btn = slide.shapes.add_shape(9, PPTXInches(bx), PPTXInches(7.02), PPTXInches(0.34), PPTXInches(0.34))
        btn.fill.solid(); btn.fill.fore_color.rgb = ACC
        btn.line.fill.background()
        tf = btn.text_frame
        tf.margin_top = 0; tf.margin_bottom = 0; tf.margin_left = 0; tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = PPTXPt(12); r.font.bold = True
        r.font.color.rgb = _rgb(pal["white"])
        try:
            btn.click_action.target_slide = prs.slides[tgt]
        except Exception:
            pass
        bx += 0.42


def _vary_pptx_layouts(slides: list, seed: int) -> list:
    """Assign a layout to each slide from its *compatible* options, shuffled per deck,
    so the same content never renders the same way twice and no two neighbours match."""
    rng = random.Random(seed)

    def compatible(sd: dict) -> list:
        opts = []
        if len(_safe_chart_data(sd.get("chart_data") or [])) >= 2:
            opts.append("chart")
        if sd.get("stats"):
            opts.append("stat")
        if sd.get("left_points") and sd.get("right_points"):
            opts.append("two_col")
        if sd.get("quote") or sd.get("callout"):
            opts.append("quote")
        if sd.get("body") and sd.get("highlights"):
            opts.append("image_text")
        if sd.get("bullets"):
            opts.append("bullets")
        if not opts:
            opts = ["stat", "image_text", "two_col"]
        return opts

    prev = None
    for i, sd in enumerate(slides):
        t = str(sd.get("type", "")).lower()
        if i == 0 or t == "title":
            sd["_layout"] = "title"; prev = "title"; continue
        if i == len(slides) - 1 or t == "closing":
            sd["_layout"] = "closing"; prev = "closing"; continue
        opts = compatible(sd)
        rng.shuffle(opts)
        choice = next((o for o in opts if o != prev), opts[0])
        sd["_layout"] = choice
        prev = choice
    return slides


def _add_native_pptx_chart(slide, chart_data, l, t, w, h, pal, kind="column"):
    """Add a real, editable PowerPoint chart styled to the deck palette. Returns the frame or None."""
    pts = _safe_chart_data(chart_data, limit=8)
    if len(pts) < 2:
        return None
    cd = CategoryChartData()
    cd.categories = [str(lbl)[:14] for lbl, _ in pts]
    cd.add_series("Series 1", [v for _, v in pts])
    ctype = XL_CHART_TYPE.LINE_MARKERS if kind == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    frame = slide.shapes.add_chart(ctype, PPTXInches(l), PPTXInches(t), PPTXInches(w), PPTXInches(h), cd)
    chart = frame.chart
    ACC = _rgb(pal["accent"]); LIGHT = _rgb(pal["light"]); WHITE = _rgb(pal["white"])
    body_font = pal.get("data_font", "Calibri")
    try:
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.gap_width = 60
        plot.has_data_labels = True
        plot.data_labels.font.size = PPTXPt(10)
        plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = WHITE
        plot.data_labels.font.name = body_font
        series = plot.series[0]
        if kind == "line":
            series.format.line.color.rgb = ACC
            series.format.line.width = PPTXPt(2.5)
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = ACC
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = PPTXPt(10)
        cat_axis.tick_labels.font.color.rgb = LIGHT
        cat_axis.tick_labels.font.name = body_font
        cat_axis.format.line.color.rgb = LIGHT
        val_axis = chart.value_axis
        val_axis.visible = False
        val_axis.has_major_gridlines = False
    except Exception:
        pass  # styling is best-effort; the chart itself still renders
    return frame


@app.post("/docs/generate/pptx")
async def generate_pptx(payload: dict):
    prompt = _normalize_topic_prompt(payload["prompt"])
    slide_count = requested_count(prompt, "slide", default=10, minimum=8, maximum=12)
    structure = await call_ai_pptx(prompt, slide_count)

    prs = Presentation()
    prs.slide_width  = PPTXInches(13.33)
    prs.slide_height = PPTXInches(7.5)

    if structure.get("design_system"):
        pal = _pptx_palette_from_structure(structure)
    else:
        pname = str(structure.get("palette", "midnight")).lower()
        pal = PPTX_PALETTES.get(pname, PPTX_PALETTES["midnight"])
        pal = {**pal, "h1_font": "Calibri", "h2_font": "Calibri", "body_font": "Calibri", "data_font": "Calibri"}

    display_font = pal.get("h1_font", "Calibri")
    heading_font = pal.get("h2_font", display_font)
    body_font = pal.get("body_font", "Calibri")
    data_font = pal.get("data_font", body_font)

    aesthetic = structure.get("aesthetic_direction") or {}
    aesthetic_label = str(aesthetic.get("label") or aesthetic.get("name") or "")

    BG    = _rgb(pal["bg"])
    PRI   = _rgb(pal["primary"])
    ACC   = _rgb(pal["accent"])
    LIGHT = _rgb(pal["light"])
    WHITE = _rgb(pal["white"])
    GRAY  = _rgb(pal["gray"])
    CARD  = _rgb(pal["card"])

    deck_title = structure.get("title", "Presentation")

    # Per-deck motion identity — chosen randomly so no two generations look alike.
    deck_seed = random.randrange(1 << 30)
    motion_rng = random.Random(deck_seed)
    deck_transition = motion_rng.choice(_PPTX_TRANSITIONS)
    entrance_keys = list(_PPTX_ENTRANCES.keys())
    motion_rng.shuffle(entrance_keys)

    # Randomised, content-aware layout assignment (different every run, no repeats).
    deck_slides = structure["slides"][:slide_count]
    _vary_pptx_layouts(deck_slides, deck_seed)

    for si, slide_data in enumerate(deck_slides):
        raw_type = str(slide_data.get("type", "stat")).lower()
        stype  = slide_data.get("_layout") or PPTX_TYPE_TO_LAYOUT.get(raw_type, raw_type)
        stitle = str(slide_data.get("title", ""))
        slide  = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # palette tuples + derived shades for gradients
        BGT, CARDT, ACCT, LIGHTT = pal["bg"], pal["card"], pal["accent"], pal["light"]
        PANEL = (0x0F, 0x12, 0x18)  # near-black data panel
        LIGHT_BG1, LIGHT_BG2 = (0xFF, 0xFF, 0xFF), (0xEE, 0xF1, 0xF7)
        INK = (0x1B, 0x24, 0x33)

        def dark_bg():
            _add_grad_rect(slide, 0, 0, 13.33, 7.5, _shade(BGT, 0.66), BGT, 125)
            _add_rect(slide, 0, 0, 0.16, 7.5, ACC)

        def light_bg():
            _add_grad_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG1, LIGHT_BG2, 125)

        def dark_footer():
            _add_rect(slide, 0, 7.24, 13.33, 0.02, ACC)
            _txt(slide, deck_title.upper(), 0.55, 7.02, 9, 0.3, size=8, color=GRAY, font=data_font)
            _txt(slide, f"{si+1:02d} / {slide_count:02d}", 11.9, 7.02, 0.9, 0.3, size=8, color=ACC, align=PP_ALIGN.RIGHT, font=data_font)

        def light_footer():
            _add_rect(slide, 0.55, 7.16, 12.23, 0.014, _rgb(_shade(LIGHTT, 0.9)))
            _txt(slide, deck_title.upper(), 0.55, 7.0, 9, 0.3, size=8, color=GRAY, font=data_font)
            _txt(slide, f"{si+1:02d} / {slide_count:02d}", 11.9, 7.0, 0.9, 0.3, size=8, color=ACC, align=PP_ALIGN.RIGHT, font=data_font)

        def eyebrow(text):
            _txt(slide, f"{si+1:02d}   {text}", 0.6, 0.5, 12, 0.32, size=9.5, bold=True, color=ACC, font=data_font)

        section_lbl = raw_type.replace("_", " ").upper()[:22]

        if stype == "title":
            _add_grad_rect(slide, 0, 0, 13.33, 7.5, _shade(BGT, 0.6), BGT, 120)
            _add_rect(slide, 0, 0, 0.16, 7.5, ACC)
            ring = slide.shapes.add_shape(9, PPTXInches(9.4), PPTXInches(-1.8), PPTXInches(6.2), PPTXInches(6.2))
            _grad(ring, _shade(CARDT, 1.12), CARDT, 120); _soft_shadow(ring, blur_pt=24, dist_pt=0, alpha_pct=22)
            ring2 = slide.shapes.add_shape(9, PPTXInches(11.1), PPTXInches(4.7), PPTXInches(3.4), PPTXInches(3.4))
            _grad(ring2, ACCT, _shade(ACCT, 0.7), 120); _soft_shadow(ring2, blur_pt=18, dist_pt=0, alpha_pct=26)

            subtitle_txt = str(slide_data.get("subtitle", ""))
            tagline_txt  = str(slide_data.get("tagline", ""))
            if aesthetic_label and aesthetic_label not in tagline_txt:
                tagline_txt = f"{aesthetic_label}  |  {tagline_txt}".strip(" |")
            if aesthetic_label:
                _txt(slide, aesthetic_label.upper(), 0.62, 1.0, 9.5, 0.4, size=10, bold=True, color=ACC, font=data_font)
            _txt(slide, stitle, 0.58, 1.72, 9.6, 3.1, size=47, bold=True, color=WHITE, font=display_font)
            if subtitle_txt:
                _txt(slide, subtitle_txt, 0.62, 4.95, 9.2, 1.1, size=20, color=LIGHT, font=body_font)
            _add_rect(slide, 0.62, 6.55, 0.9, 0.05, ACC)
            if tagline_txt:
                _txt(slide, tagline_txt, 0.62, 6.72, 11, 0.5, size=11.5, color=GRAY, font=body_font)
            _txt(slide, f"01 / {slide_count:02d}", 11.9, 6.72, 0.9, 0.4, size=11, color=GRAY, align=PP_ALIGN.RIGHT, font=data_font)

        elif stype == "closing":
            _add_grad_rect(slide, 0, 0, 13.33, 7.5, _shade(BGT, 0.6), BGT, 120)
            _add_rect(slide, 0, 0, 0.16, 7.5, ACC)
            ring = slide.shapes.add_shape(9, PPTXInches(8.2), PPTXInches(0.4), PPTXInches(6.2), PPTXInches(6.2))
            _grad(ring, _shade(CARDT, 1.1), CARDT, 120); _soft_shadow(ring, blur_pt=24, dist_pt=0, alpha_pct=22)
            cta = str(slide_data.get("cta", "Thank you."))
            tagline = str(slide_data.get("tagline", ""))
            _txt(slide, "IN CLOSING", 0.62, 0.7, 5, 0.4, size=10, color=ACC, bold=True, font=data_font)
            _txt(slide, stitle, 0.58, 1.5, 9.4, 2.6, size=40, bold=True, color=WHITE, font=display_font)
            cta_box = _round_rect(slide, 0.6, 4.35, 7.6, 1.05, grad=(ACCT, _shade(ACCT, 0.82), 90), shadow=True, radius=0.16)
            ctf = cta_box.text_frame; ctf.word_wrap = True
            ctf.paragraphs[0].text = cta
            ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
            r0 = ctf.paragraphs[0].runs[0]
            r0.font.size = PPTXPt(17); r0.font.bold = True
            r0.font.color.rgb = _rgb(_shade(BGT, 0.8)); r0.font.name = body_font
            if tagline:
                _txt(slide, tagline, 0.62, 5.65, 9, 0.6, size=13.5, color=LIGHT, italic=True, font=body_font)
            dark_footer()

        elif stype == "bullets":
            light_bg()
            eyebrow(stitle.upper()[:44])
            _txt(slide, stitle, 0.58, 0.92, 12, 1.1, size=30, bold=True, color=_rgb(BGT), font=heading_font)
            _add_rect(slide, 0.62, 1.95, 0.9, 0.05, ACC)
            bullets = slide_data.get("bullets", [])
            _add_bullets_textbox(slide, bullets, 0.6, 2.35, 12.1, 4.5, size=15, color=_rgb(INK), marker="▸", font=body_font)
            light_footer()

        elif stype == "two_col":
            light_bg()
            eyebrow(section_lbl or "COMPARISON")
            _txt(slide, stitle, 0.58, 0.92, 12.2, 1.0, size=27, bold=True, color=_rgb(BGT), font=heading_font)
            _round_rect(slide, 0.55, 2.0, 6.05, 4.85, grad=(LIGHT_BG1, _shade(LIGHTT, 1.06), 120), shadow=True, radius=0.045)
            _add_rect(slide, 0.9, 2.35, 1.5, 0.06, ACC)
            _txt(slide, str(slide_data.get("left_title", "") or "Perspective A").upper(), 0.9, 2.5, 5.2, 0.4, size=10, bold=True, color=ACC, font=data_font)
            _add_bullets_textbox(slide, slide_data.get("left_points", []), 0.9, 3.05, 5.4, 3.6, size=13.5, color=_rgb(INK), marker="→", font=body_font)
            _round_rect(slide, 6.9, 2.0, 5.9, 4.85, grad=(_shade(CARDT, 1.08), CARDT, 120), shadow=True, radius=0.045)
            _add_rect(slide, 7.25, 2.35, 1.5, 0.06, ACC)
            _txt(slide, str(slide_data.get("right_title", "") or "Perspective B").upper(), 7.25, 2.5, 5.2, 0.4, size=10, bold=True, color=ACC, font=data_font)
            _add_bullets_textbox(slide, slide_data.get("right_points", []), 7.25, 3.05, 5.2, 3.6, size=13.5, color=WHITE, marker="→", font=body_font)
            light_footer()

        elif stype == "stat":
            dark_bg()
            eyebrow(section_lbl or "KEY METRICS")
            _txt(slide, stitle, 0.58, 0.95, 12.2, 1.1, size=30, bold=True, color=WHITE, font=heading_font)
            _add_rect(slide, 0.62, 2.05, 0.9, 0.05, ACC)
            stats_data = slide_data.get("stats", [])[:3]
            card_w, gap = 3.92, 0.28
            starts = [0.6 + i * (card_w + gap) for i in range(3)]
            for ci, st in enumerate(stats_data):
                cx = starts[ci]
                _round_rect(slide, cx, 2.5, card_w, 3.95, grad=(_shade(CARDT, 1.12), CARDT, 120), shadow=True, radius=0.05)
                _add_rect(slide, cx + 0.35, 2.5, 1.0, 0.07, ACC)
                _txt(slide, str(st.get("number", "—")), cx + 0.25, 3.15, card_w - 0.5, 1.5, size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=data_font)
                _txt(slide, str(st.get("label", "")), cx + 0.25, 4.75, card_w - 0.5, 0.55, size=12.5, bold=True, color=ACC, align=PP_ALIGN.CENTER, font=heading_font)
                _txt(slide, str(st.get("context", "")), cx + 0.3, 5.3, card_w - 0.6, 1.0, size=10.5, color=LIGHT, align=PP_ALIGN.CENTER, font=body_font)
            dark_footer()

        elif stype == "quote":
            _add_grad_rect(slide, 0, 0, 13.33, 7.5, _shade(BGT, 0.62), BGT, 125)
            _add_rect(slide, 0, 0, 0.16, 7.5, ACC)
            _txt(slide, "“", 0.45, -0.35, 3, 2.6, size=150, color=_rgb(_shade(CARDT, 1.2)), bold=True, font=display_font)
            quote_txt  = str(slide_data.get("quote", "") or slide_data.get("callout", ""))
            attrib_txt = str(slide_data.get("attribution", ""))
            _txt(slide, quote_txt, 1.0, 1.75, 8.9, 3.6, size=25, color=WHITE, italic=True, font=display_font)
            _add_rect(slide, 1.05, 5.5, 0.8, 0.05, ACC)
            if attrib_txt:
                _txt(slide, attrib_txt, 1.05, 5.7, 8.5, 0.5, size=13, color=ACC, bold=True, font=data_font)
            hi = [h for h in (slide_data.get("highlights", []) or slide_data.get("bullets", [])) if str(h).strip()]
            for hj, h in enumerate(hi[:3]):
                hy = 1.35 + hj * 1.75
                _round_rect(slide, 10.35, hy, 2.5, 1.5, grad=(_shade(CARDT, 1.1), CARDT, 120), shadow=True, radius=0.08)
                _add_rect(slide, 10.55, hy + 0.22, 0.45, 0.05, ACC)
                _txt(slide, str(h), 10.55, hy + 0.38, 2.1, 1.0, size=10.5, color=WHITE, font=body_font)
            dark_footer()

        elif stype == "chart":
            dark_bg()
            eyebrow(section_lbl or "DATA")
            _txt(slide, stitle, 0.58, 0.95, 12.2, 1.0, size=29, bold=True, color=WHITE, font=heading_font)
            _add_rect(slide, 0.62, 2.0, 0.9, 0.05, ACC)
            _round_rect(slide, 0.55, 2.4, 8.35, 4.35, grad=(_shade(PANEL, 1.5), PANEL, 120), shadow=True, radius=0.04)
            chart_kind = "line" if ("trend" in (stitle + raw_type).lower() or "growth" in (stitle + raw_type).lower()) else "column"
            frame = _add_native_pptx_chart(slide, slide_data.get("chart_data"), 0.85, 2.75, 7.75, 3.7, pal, kind=chart_kind)
            if frame is None:
                stats_data = slide_data.get("stats", [])[:3]
                for ci, st in enumerate(stats_data):
                    cx = [0.85, 3.55, 6.25][ci]
                    _txt(slide, str(st.get("number", "—")), cx, 3.4, 2.5, 1.2, size=34, bold=True, color=WHITE, font=data_font)
                    _txt(slide, str(st.get("label", "")), cx, 4.6, 2.5, 0.6, size=11, bold=True, color=ACC, font=heading_font)
            _round_rect(slide, 9.15, 2.4, 3.63, 4.35, grad=(_shade(CARDT, 1.1), CARDT, 120), shadow=True, radius=0.05)
            _add_rect(slide, 9.4, 2.75, 1.0, 0.06, ACC)
            _txt(slide, "WHAT IT MEANS", 9.4, 2.9, 3.2, 0.4, size=10, bold=True, color=ACC, font=data_font)
            rail = slide_data.get("highlights") or [st.get("context", "") for st in slide_data.get("stats", [])]
            _add_bullets_textbox(slide, [r for r in rail if r][:4], 9.4, 3.45, 3.2, 3.1, size=11.5, color=WHITE, marker="◆", font=body_font)
            dark_footer()

        elif stype == "image_text":
            light_bg()
            eyebrow(section_lbl or "OVERVIEW")
            _txt(slide, stitle, 0.58, 0.92, 8.6, 1.0, size=27, bold=True, color=_rgb(BGT), font=heading_font)
            _add_rect(slide, 0.62, 1.95, 0.9, 0.05, ACC)
            body_txt = str(slide_data.get("body", ""))
            _txt(slide, body_txt, 0.6, 2.35, 8.1, 4.4, size=14, color=_rgb(INK), font=body_font)
            _round_rect(slide, 9.0, 2.0, 3.78, 4.85, grad=(_shade(CARDT, 1.1), CARDT, 120), shadow=True, radius=0.05)
            _add_rect(slide, 9.3, 2.35, 1.0, 0.06, ACC)
            _txt(slide, "KEY DATA", 9.3, 2.5, 3.3, 0.4, size=10, bold=True, color=ACC, font=data_font)
            _add_bullets_textbox(slide, slide_data.get("highlights", [])[:4], 9.3, 3.05, 3.3, 3.6, size=12, color=WHITE, marker="◆", font=body_font)
            light_footer()

        else:
            light_bg()
            eyebrow(section_lbl or "INSIGHT")
            _txt(slide, stitle, 0.58, 0.92, 12.2, 1.0, size=27, bold=True, color=_rgb(BGT), font=heading_font)
            _add_rect(slide, 0.62, 1.95, 0.9, 0.05, ACC)
            stats_fb = slide_data.get("stats", [])[:3]
            if stats_fb:
                card_w, gap = 3.92, 0.28
                starts = [0.6 + i * (card_w + gap) for i in range(3)]
                for ci, st in enumerate(stats_fb):
                    cx = starts[ci]
                    _round_rect(slide, cx, 2.5, card_w, 3.9, grad=(_shade(CARDT, 1.12), CARDT, 120), shadow=True, radius=0.05)
                    _add_rect(slide, cx + 0.35, 2.5, 1.0, 0.07, ACC)
                    _txt(slide, str(st.get("number", "—")), cx + 0.25, 3.2, card_w - 0.5, 1.4, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=data_font)
                    _txt(slide, str(st.get("label", "")), cx + 0.25, 4.75, card_w - 0.5, 0.6, size=12, bold=True, color=ACC, align=PP_ALIGN.CENTER, font=heading_font)
            else:
                _add_bullets_textbox(slide, slide_data.get("bullets", []), 0.6, 2.35, 12.1, 4.4, size=14, color=_rgb(INK), marker="▸", font=body_font)
            light_footer()

        # ── motion: cohesive deck transition + auto-playing staggered entrance ──
        _set_slide_transition(slide, deck_transition)
        _animate_slide(slide, entrance_keys[si % len(entrance_keys)])

    # Clickable navigation — added after all slides exist so jump targets resolve.
    all_slides = list(prs.slides)
    for ni, slide in enumerate(all_slides):
        _add_nav_controls(slide, prs, ni, len(all_slides), pal)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    blob = await upload_to_vercel_blob(
        buffer, "output.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    return blob


# ══════════════════════════════════════════════════════════════════════════════
# DOCX GENERATION
# ══════════════════════════════════════════════════════════════════════════════

def _parse_output_format(prompt: str) -> str:
    """Detect desired output from prompt text. Returns pdf|pptx|docx|both."""
    normalized = prompt.lower()
    if re.search(r"\b(?:both|pdf\s*\+\s*pptx|pdf\s*and\s*pptx|all\s+formats)\b", normalized):
        return "both"
    if re.search(r"\b(?:pptx|powerpoint|presentation|slides?)\b", normalized):
        return "pptx"
    if re.search(r"\b(?:docx|word|document)\b", normalized) and "pdf" not in normalized:
        return "docx"
    if re.search(r"\b(?:pdf|report)\b", normalized):
        return "pdf"
    return "pdf"


def _normalize_topic_prompt(prompt: str) -> str:
    """Strip template placeholders from user prompt."""
    cleaned = prompt
    for placeholder in (
        r"\[INSERT TOPIC\]", r"\[INSERT TOPIC HERE\]",
        r"\[PDF\s*/\s*PPTX\s*/\s*DOCS/BOTH\]", r"\[PDF\s*/\s*PPTX\]",
        r"Topic:\s*\[.*?\]", r"Output format:\s*\[.*?\]",
    ):
        cleaned = re.sub(placeholder, "", cleaned, flags=re.IGNORECASE)
    return cleaned.strip() or prompt.strip()


@app.post("/docs/generate")
async def generate_document(payload: dict):
    """Unified endpoint — routes to PDF, PPTX, DOCX, or both based on prompt."""
    prompt = _normalize_topic_prompt(payload.get("prompt", ""))
    fmt = payload.get("format") or _parse_output_format(prompt)
    fmt = str(fmt).lower().strip()

    if fmt == "both":
        pdf_result = await generate_pdf({"prompt": prompt})
        pptx_result = await generate_pptx({"prompt": prompt})
        return {"pdf": pdf_result, "pptx": pptx_result}
    if fmt in ("pptx", "presentation", "slides"):
        return await generate_pptx({"prompt": prompt})
    if fmt in ("docx", "word", "document"):
        return await generate_docx({"prompt": prompt})
    return await generate_pdf({"prompt": prompt})


async def call_ai_docx(prompt: str) -> dict:
    core = _creative_director_core("section", 6, "section")
    system = f"""{core}

Create a polished Word document. Return ONLY valid JSON, no markdown fences.

Adapt the creative director principles to a flowing document (not slides):
- Include aesthetic_direction and design_system (5 hex colors, typography)
- Use REAL topic-specific statistics woven into prose
- NO bullet-list-heavy sections — prefer narrative paragraphs with embedded data
- Include pull-quotes and a data_table in at least one section

JSON schema:
{{
  "title": "Document Title",
  "subtitle": "Compelling subtitle",
  "author": "Prepared by AI Document Service",
  "date": "2026",
  "abstract": "2-3 sentence executive summary with at least one real statistic.",
  "aesthetic_direction": {{"name": "luxury_editorial", "label": "Luxury Editorial", "rationale": "..."}},
  "design_system": {{
    "colors": {{"dominant": "#111", "secondary": "#333", "supporting_1": "#666", "supporting_2": "#EEE", "accent": "#C9A84C"}},
    "typography": {{"h1_font": "Cormorant Garamond", "body_font": "Libre Franklin"}}
  }},
  "palette": "indigo|teal|crimson|emerald|slate|violet|amber|rose",
  "sections": [
    {{
      "heading": "Section Heading",
      "body": "Rich detailed paragraph, 80+ words with real figures and dates.",
      "callout": "Optional pull-quote for this section.",
      "data_table": {{"headers": ["Metric", "Value"], "rows": [["Example", "42%"]]}},
      "subsections": [{{"heading": "Subsection", "body": "40+ words."}}]
    }}
  ],
  "conclusion": "Compelling 2-3 sentence conclusion with forward-looking tone."
}}

RULES:
- 6-8 sections minimum
- Every section body: 80+ words, topic-specific
- palette: pick one that emotionally fits the topic"""

    raw = await _deepinfra_call([
        {"role": "system", "content": system},
        {"role": "user", "content": prompt},
    ], temperature=0.38)
    return _as_json(raw, "DOCX content generator")


def _set_cell_bg(cell, hex_color: str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color.lstrip('#'))
    tcPr.append(shd)


def _set_outline_level(paragraph, level: int):
    """Tag a paragraph with a Word outline level so it appears in the Navigation pane
    and is picked up by a TOC field (\\u switch) — independent of paragraph style."""
    pPr = paragraph._p.get_or_add_pPr()
    existing = pPr.find(qn('w:outlineLvl'))
    if existing is not None:
        pPr.remove(existing)
    ol = OxmlElement('w:outlineLvl')
    ol.set(qn('w:val'), str(level))
    pPr.append(ol)


def _add_docx_toc(doc: Document):
    """Insert a live, clickable Table of Contents field (updates on open / right-click)."""
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run()
    r = run._r
    begin = OxmlElement('w:fldChar'); begin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText'); instr.set(qn('xml:space'), 'preserve')
    instr.text = 'TOC \\o "1-3" \\h \\z \\u'
    sep = OxmlElement('w:fldChar'); sep.set(qn('w:fldCharType'), 'separate')
    placeholder = OxmlElement('w:t')
    placeholder.text = "Table of contents — right-click and choose “Update Field”."
    end = OxmlElement('w:fldChar'); end.set(qn('w:fldCharType'), 'end')
    for child in (begin, instr, sep, placeholder, end):
        r.append(child)


def _enable_update_fields(doc: Document):
    """Ask Word to refresh fields (the TOC) automatically when the document opens."""
    settings = doc.settings.element
    if settings.find(qn('w:updateFields')) is None:
        el = OxmlElement('w:updateFields')
        el.set(qn('w:val'), 'true')
        settings.append(el)


def _add_docx_heading_band(doc: Document, text: str, bg_hex: str, fg_hex: str, level_pt: int = 14,
                            outline_level: int = None, font: str = None):
    table = doc.add_table(rows=1, cols=1)
    table.style = 'Table Grid'
    cell = table.cell(0, 0)
    cell.width = Inches(6.5)
    _set_cell_bg(cell, bg_hex)
    tbl = table._tbl
    tblPr = tbl.find(qn('w:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    tblBorders = OxmlElement('w:tblBorders')
    for side in ['top','left','bottom','right','insideH','insideV']:
        border = OxmlElement(f'w:{side}')
        border.set(qn('w:val'), 'none')
        tblBorders.append(border)
    tblPr.append(tblBorders)
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(7)
    p.paragraph_format.space_after  = Pt(7)
    p.paragraph_format.left_indent  = Inches(0.16)
    run = p.add_run(text)
    _docx_run(run, font=font, size=level_pt, bold=True,
              color=RGBColor(int(fg_hex[1:3],16), int(fg_hex[3:5],16), int(fg_hex[5:7],16)))
    if outline_level is not None:
        _set_outline_level(p, outline_level)
    doc.add_paragraph()


def _add_pull_quote(doc: Document, text: str, accent_hex: str):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent  = Inches(0.4)
    p.paragraph_format.right_indent = Inches(0.4)
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after  = Pt(8)
    run = p.add_run(f'"{text}"')
    run.italic = True
    run.font.size = Pt(12)
    r,g,b = int(accent_hex[1:3],16),int(accent_hex[3:5],16),int(accent_hex[5:7],16)
    run.font.color.rgb = RGBColor(r,g,b)
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    left = OxmlElement('w:left')
    left.set(qn('w:val'), 'single')
    left.set(qn('w:sz'), '24')
    left.set(qn('w:space'), '12')
    left.set(qn('w:color'), accent_hex.lstrip('#'))
    pBdr.append(left)
    pPr.append(pBdr)


def _add_docx_data_table(doc: Document, table: dict, header_hex: str, accent_hex: str, alt_hex: str):
    """Render a clean, banded data table from {headers:[...], rows:[[...]]}."""
    headers = (table or {}).get("headers") or []
    rows = (table or {}).get("rows") or []
    if not headers or not rows:
        return
    headers = [str(h) for h in headers[:6]]
    n = len(headers)
    t = doc.add_table(rows=1, cols=n)
    t.style = 'Table Grid'
    t.autofit = True
    hdr_cells = t.rows[0].cells
    for ci, h in enumerate(headers):
        _set_cell_bg(hdr_cells[ci], header_hex.lstrip('#'))
        para = hdr_cells[ci].paragraphs[0]
        para.paragraph_format.space_before = Pt(2)
        para.paragraph_format.space_after = Pt(2)
        run = para.add_run(h)
        run.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for ri, row in enumerate(rows[:12]):
        cells = t.add_row().cells
        vals = (list(row) + [""] * n)[:n]
        for ci, v in enumerate(vals):
            if ri % 2 == 1:
                _set_cell_bg(cells[ci], alt_hex.lstrip('#'))
            para = cells[ci].paragraphs[0]
            para.paragraph_format.space_before = Pt(2)
            para.paragraph_format.space_after = Pt(2)
            run = para.add_run(str(v))
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x1F, 0x2A, 0x3C)
            if ci == 0:
                run.bold = True
                run.font.color.rgb = RGBColor(
                    int(accent_hex[1:3], 16), int(accent_hex[3:5], 16), int(accent_hex[5:7], 16))
    doc.add_paragraph()


def _docx_hexrgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _docx_set_base_fonts(doc: Document, body_font: str, body_pt: float = 10.5,
                          body_color: str = "26303F"):
    """Set the document's default (Normal) typeface, size, and colour so every
    paragraph inherits the AI-chosen body font instead of Calibri."""
    style = doc.styles['Normal']
    style.font.name = body_font
    style.font.size = Pt(body_pt)
    style.font.color.rgb = _docx_hexrgb(body_color)
    # Ensure the default run properties carry the font for all script ranges.
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.find(qn('w:rFonts'))
    if rfonts is None:
        rfonts = OxmlElement('w:rFonts'); rpr.append(rfonts)
    for attr in ('w:ascii', 'w:hAnsi', 'w:cs'):
        rfonts.set(qn(attr), body_font)
    style.paragraph_format.line_spacing = 1.35
    style.paragraph_format.space_after = Pt(6)


def _docx_run(run, font=None, size=None, color=None, bold=None, italic=None,
              caps=False, spacing=None):
    if font:
        run.font.name = font
        rpr = run._r.get_or_add_rPr()
        rfonts = rpr.find(qn('w:rFonts'))
        if rfonts is None:
            rfonts = OxmlElement('w:rFonts'); rpr.append(rfonts)
        for attr in ('w:ascii', 'w:hAnsi', 'w:cs'):
            rfonts.set(qn(attr), font)
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = _docx_hexrgb(color) if isinstance(color, str) else color
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if caps:
        run.font.all_caps = True
    if spacing is not None:  # letter-spacing in twips
        rpr = run._r.get_or_add_rPr()
        sp = OxmlElement('w:spacing'); sp.set(qn('w:val'), str(spacing)); rpr.append(sp)
    return run


def _docx_accent_bar(doc: Document, hex_color: str, width_dxa: int = 1300, height_pt: float = 3.2,
                     space_before: int = 6, space_after: int = 6):
    """A short, left-aligned solid accent rule (a borderless 1-cell shaded table)."""
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    _set_cell_bg(cell, hex_color.lstrip('#'))
    tbl = table._tbl
    tblPr = tbl.find(qn('w:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr'); tbl.insert(0, tblPr)
    borders = OxmlElement('w:tblBorders')
    for side in ('top', 'left', 'bottom', 'right', 'insideH', 'insideV'):
        b = OxmlElement(f'w:{side}'); b.set(qn('w:val'), 'none'); borders.append(b)
    tblPr.append(borders)
    w = tblPr.find(qn('w:tblW'))
    if w is not None:
        tblPr.remove(w)
    tblW = OxmlElement('w:tblW'); tblW.set(qn('w:w'), str(width_dxa)); tblW.set(qn('w:type'), 'dxa')
    tblPr.append(tblW)
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    r = p.add_run(); r.font.size = Pt(height_pt)
    return table


def _docx_page_footer(doc: Document, left_text: str, accent_hex: str, gray_hex: str, data_font: str):
    """Running footer: document title on the left, live page number on the right."""
    section = doc.sections[0]
    footer = section.footer
    footer.is_linked_to_previous = False
    p = footer.paragraphs[0]
    p.text = ""
    # right tab at the content width so the page number sits flush-right
    pPr = p._p.get_or_add_pPr()
    tabs = OxmlElement('w:tabs')
    tab = OxmlElement('w:tab'); tab.set(qn('w:val'), 'right')
    content_w = int((section.page_width - section.left_margin - section.right_margin) / 635)  # EMU→twips
    tab.set(qn('w:pos'), str(content_w)); tabs.append(tab); pPr.append(tabs)
    lr = p.add_run(left_text.upper()); _docx_run(lr, font=data_font, size=7.5, color=gray_hex, spacing=20)
    p.add_run("\t")
    pr = p.add_run("PAGE "); _docx_run(pr, font=data_font, size=7.5, color=accent_hex.lstrip('#'))
    fld = p.add_run()
    begin = OxmlElement('w:fldChar'); begin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText'); instr.set(qn('xml:space'), 'preserve'); instr.text = 'PAGE'
    end = OxmlElement('w:fldChar'); end.set(qn('w:fldCharType'), 'end')
    fld._r.append(begin); fld._r.append(instr); fld._r.append(end)
    _docx_run(fld, font=data_font, size=7.5, color=accent_hex.lstrip('#'), bold=True)


@app.post("/docs/generate/docx")
async def generate_docx(payload: dict):
    structure = await call_ai_docx(_normalize_topic_prompt(payload["prompt"]))

    pname = str(structure.get("palette", "indigo")).lower()
    pal = PALETTES.get(pname, PALETTES["indigo"])
    tokens = _resolve_design_tokens(structure)
    PRI_HEX = tokens["primary"]
    ACC_HEX = tokens["accent"]
    pal = {**pal, "light": tokens["light"]}
    H1_FONT = tokens["h1_font"]
    H2_FONT = tokens["h2_font"]
    BODY_FONT = tokens["body_font"]
    DATA_FONT = tokens["data_font"]

    def hexrgb(h):
        h = h.lstrip('#')
        return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

    PRI_RGB  = hexrgb(PRI_HEX)
    ACC_RGB  = hexrgb(ACC_HEX)
    GRAY_RGB = RGBColor(0x6B,0x72,0x80)
    BODY_RGB = RGBColor(0x26,0x30,0x3F)

    doc = Document()
    _docx_set_base_fonts(doc, BODY_FONT, body_pt=10.5)

    for sec in doc.sections:
        sec.top_margin    = Inches(1.15)
        sec.bottom_margin = Inches(1.0)
        sec.left_margin   = Inches(1.15)
        sec.right_margin  = Inches(1.15)

    # ── editorial cover ──────────────────────────────────────────────────────
    aesthetic = structure.get("aesthetic_direction") or {}
    eyebrow_txt = str(aesthetic.get("label") or aesthetic.get("name") or "Report").upper()
    spacer = doc.add_paragraph(); spacer.paragraph_format.space_before = Pt(90)

    eyebrow_para = doc.add_paragraph()
    eyebrow_para.paragraph_format.space_after = Pt(10)
    _docx_run(eyebrow_para.add_run(eyebrow_txt), font=DATA_FONT, size=10.5, color=ACC_HEX.lstrip('#'),
              bold=True, spacing=48)

    title_para = doc.add_paragraph()
    title_para.paragraph_format.space_after = Pt(12)
    title_para.paragraph_format.line_spacing = 1.0
    _docx_run(title_para.add_run(structure.get("title", "Document")), font=H1_FONT, size=38,
              color=PRI_HEX.lstrip('#'), bold=True)

    if structure.get("subtitle"):
        sp = doc.add_paragraph()
        sp.paragraph_format.space_after = Pt(18)
        _docx_run(sp.add_run(structure["subtitle"]), font=H2_FONT, size=15,
                  color="55606E", italic=False)

    _docx_accent_bar(doc, ACC_HEX, width_dxa=1500, height_pt=3.4, space_before=2, space_after=16)

    meta_para = doc.add_paragraph()
    meta_para.paragraph_format.space_after = Pt(3)
    author_txt = str(structure.get("author", "")).strip()
    date_txt = str(structure.get("date", "")).strip()
    meta_line = "  ·  ".join([t for t in (author_txt, date_txt) if t])
    _docx_run(meta_para.add_run(meta_line), font=DATA_FONT, size=9.5, color="6B7280", spacing=20)

    doc.add_page_break()

    # Clickable, auto-updating table of contents (driven by paragraph outline levels).
    if structure.get("sections"):
        _add_docx_heading_band(doc, "Contents", PRI_HEX, "#FFFFFF", level_pt=13, font=H2_FONT)
        _add_docx_toc(doc)
        doc.add_page_break()

    if structure.get("abstract"):
        _add_docx_heading_band(doc, "Executive Summary", PRI_HEX, "#FFFFFF", level_pt=13,
                               outline_level=0, font=H2_FONT)
        abs_para = doc.add_paragraph(structure["abstract"])
        abs_para.paragraph_format.space_after = Pt(12)
        abs_para.paragraph_format.left_indent = Inches(0.28)
        abs_para.paragraph_format.line_spacing = 1.4
        # left accent rule on the lede
        pBdr = OxmlElement('w:pBdr'); left_b = OxmlElement('w:left')
        left_b.set(qn('w:val'), 'single'); left_b.set(qn('w:sz'), '18')
        left_b.set(qn('w:space'), '14'); left_b.set(qn('w:color'), ACC_HEX.lstrip('#'))
        pBdr.append(left_b); abs_para._p.get_or_add_pPr().append(pBdr)
        for run in abs_para.runs:
            _docx_run(run, font=H2_FONT, size=12.5, color="3A4453")
        doc.add_paragraph()

    for si, sec in enumerate(structure.get("sections", [])):
        is_alt = (si % 2 == 1)
        h_bg  = ACC_HEX if is_alt else PRI_HEX
        section_label = f"{si+1:02d}   {sec['heading']}"
        _add_docx_heading_band(doc, section_label, h_bg, "#FFFFFF", level_pt=14,
                               outline_level=0, font=H2_FONT)

        body_para = doc.add_paragraph(sec.get("body",""))
        body_para.paragraph_format.space_after  = Pt(11)
        body_para.paragraph_format.space_before = Pt(3)
        body_para.paragraph_format.line_spacing = 1.4
        body_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        for run in body_para.runs:
            run.font.size = Pt(11)
            run.font.color.rgb = BODY_RGB

        if sec.get("callout"):
            _add_pull_quote(doc, sec["callout"], ACC_HEX)

        if sec.get("data_table"):
            _add_docx_data_table(doc, sec["data_table"], PRI_HEX, ACC_HEX,
                                 pal.get("light", "#E0E7FF"))

        for sub_idx, sub in enumerate(sec.get("subsections", [])):
            sub_para = doc.add_paragraph()
            sub_para.paragraph_format.space_before = Pt(11)
            sub_para.paragraph_format.space_after  = Pt(4)
            sub_label = f"{si+1}.{sub_idx+1}   {sub['heading']}"
            _docx_run(sub_para.add_run(sub_label), font=H2_FONT, size=13, bold=True, color=ACC_HEX.lstrip('#'))
            _set_outline_level(sub_para, 1)

            sub_body = doc.add_paragraph(sub.get("body",""))
            sub_body.paragraph_format.space_after = Pt(8)
            sub_body.paragraph_format.left_indent = Inches(0.2)
            sub_body.paragraph_format.line_spacing = Pt(16)
            for run in sub_body.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = BODY_RGB

        # Horizontal separator between sections
        sep_para = doc.add_paragraph()
        sep_para.paragraph_format.space_before = Pt(6)
        sep_para.paragraph_format.space_after  = Pt(6)
        pPr_sep = sep_para._p.get_or_add_pPr()
        pBdr_sep = OxmlElement('w:pBdr')
        bottom_sep = OxmlElement('w:bottom')
        bottom_sep.set(qn('w:val'), 'single')
        bottom_sep.set(qn('w:sz'), '4')
        bottom_sep.set(qn('w:space'), '1')
        bottom_sep.set(qn('w:color'), pal.get('light', '#E0E7FF').lstrip('#'))
        pBdr_sep.append(bottom_sep)
        pPr_sep.append(pBdr_sep)

    if structure.get("conclusion"):
        _add_docx_heading_band(doc, "Conclusion", PRI_HEX, "#FFFFFF", level_pt=13,
                               outline_level=0, font=H2_FONT)
        conc_para = doc.add_paragraph(structure["conclusion"])
        conc_para.paragraph_format.space_after = Pt(10)
        conc_para.paragraph_format.line_spacing = 1.4
        for run in conc_para.runs:
            _docx_run(run, font=H2_FONT, size=11.5, color="3A4453")

    _docx_page_footer(doc, structure.get("title", "Document"), ACC_HEX, "9AA2AD", DATA_FONT)
    _enable_update_fields(doc)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    blob = await upload_to_vercel_blob(
        buffer, "output.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    return blob
