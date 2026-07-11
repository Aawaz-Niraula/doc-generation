"""render_pptx.py — 16:9 presentation renderer (python-pptx).

Consumes the same DocumentSpec as render_pdf / render_docx. Everything is
positioned explicitly in inches on blank layouts — no default placeholders.

Design-system rules enforced here:
- metric-safe Office fonts only (theme.pptx_fonts — Cambria/Calibri/Arial/
  Consolas), serif display over sans body;
- no accent bars under titles, no edge stripes, no single-edge card strips —
  depth comes from gradients, rounded cards, and soft shadows; the repeated
  motif is the numbered accent circle;
- title + closing slides on the dominant dark color, content slides light
  (sandwich) or a fully dark premium deck, chosen per theme;
- every text box is word-wrapped and passed through the auto-fit guard
  (theme.fit_text_size) so text can never overflow its container;
- speaker notes on every slide, document properties set, native editable
  charts, real photos (cover-cropped), auto-playing entrance animations and
  slide transitions injected as OOXML, clickable prev/home/next navigation.
"""

from __future__ import annotations

import io
import random
import re
from datetime import datetime, timezone
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

from pptx.opc.constants import RELATIONSHIP_TYPE as RT

import assets
from content_model import DocumentSpec, Section
from theme import (GAP, MARGIN, SLIDE_H, SLIDE_W, Theme, clip_chars, clip_words,
                   estimate_text_height_in, fit_text_size, hex_to_rgb,
                   no_consecutive_repeat, readable_on)

# ─── Color / shape primitives ─────────────────────────────────────────────────


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(hex_color))


def _shade(hex_color: str, f: float) -> str:
    r, g, b = (max(0, min(255, int(round(v * f)))) for v in hex_to_rgb(hex_color))
    return "#{:02X}{:02X}{:02X}".format(r, g, b)


def _solid(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def _grad(shape, c1: str, c2: str, angle: float = 90.0):
    """Two-stop linear gradient fill; falls back to solid on any failure."""
    try:
        f = shape.fill
        f.gradient()
        stops = f.gradient_stops
        stops[0].color.rgb = _rgb(c1)
        stops[1].color.rgb = _rgb(c2)
        try:
            f.gradient_angle = angle
        except Exception:
            pass
    except Exception:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(c1)
    shape.line.fill.background()
    return shape


def _soft_shadow(shape, blur_pt=13, dist_pt=7, dir_deg=90, alpha_pct=30, color="0B1220"):
    """Soft outer drop shadow (python-pptx has no high-level API for it)."""
    try:
        spPr = shape._element.spPr
        old = spPr.find(qn('a:effectLst'))
        if old is not None:
            spPr.remove(old)
        blur = int(blur_pt * 12700); dist = int(dist_pt * 12700)
        direction = int(dir_deg * 60000); alpha = int(alpha_pct * 1000)
        xml = (f'<a:effectLst {nsdecls("a")}>'
               f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{direction}" rotWithShape="0">'
               f'<a:srgbClr val="{color}"><a:alpha val="{alpha}"/></a:srgbClr>'
               f'</a:outerShdw></a:effectLst>')
        spPr.append(parse_xml(xml))
    except Exception:
        pass
    return shape


def _rect(slide, l, t, w, h, color: str):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    return _solid(shape, color)


def _grad_rect(slide, l, t, w, h, c1: str, c2: str, angle=90.0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    return _grad(shape, c1, c2, angle)


def _round_card(slide, l, t, w, h, color: str = None, grad: Tuple[str, str, float] = None,
                shadow: bool = True, radius: float = 0.06):
    shape = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))  # rounded rect
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    if grad is not None:
        _grad(shape, grad[0], grad[1], grad[2] if len(grad) > 2 else 90.0)
    else:
        _solid(shape, color)
    if shadow:
        _soft_shadow(shape)
    return shape


def _oval(slide, l, t, d, color: str):
    shape = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    return _solid(shape, color)


def _dot_grid(slide, l, t, cols: int, rows: int, color: str,
              dot: float = 0.055, gap: float = 0.21) -> int:
    """Editorial dot-matrix accent (native vector, crisp at any zoom).
    Returns the number of shapes added so callers can keep them static."""
    n = 0
    for r in range(rows):
        for c in range(cols):
            _oval(slide, l + c * gap, t + r * gap, dot, color)
            n += 1
    return n


def _outline_ring(slide, l, t, d_in: float, color: str, weight: float = 2.25):
    """Unfilled circle outline — a floating vector accent over photos."""
    shape = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(d_in), Inches(d_in))
    shape.fill.background()
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(weight)
    return shape


def _pct_of(value: str) -> Optional[float]:
    """Extract a 0–100 percentage from a stat value like '95%' or '37.5 %'."""
    m = re.search(r"(\d+(?:\.\d+)?)\s*%", str(value or ""))
    if not m:
        return None
    pct = float(m.group(1))
    return pct if 0 < pct <= 100 else None


def _progress_bar(slide, l, t, w, pct: float, accent: str, track: str, h: float = 0.13):
    """Vector percentage bar: full-width track + accent fill. The 'SVG gauge'
    of the deck — reads instantly and survives every Office/Keynote import."""
    track_bar = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        track_bar.adjustments[0] = 0.5
    except Exception:
        pass
    _solid(track_bar, track)
    fill_w = max(w * pct / 100.0, h)
    fill_bar = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(fill_w), Inches(h))
    try:
        fill_bar.adjustments[0] = 0.5
    except Exception:
        pass
    _solid(fill_bar, accent)


def _num_circle(slide, l, t, d, idx: int, accent: str, font: str):
    """The deck's repeated motif: a numbered accent circle."""
    circ = _oval(slide, l, t, d, accent)
    tf = circ.text_frame
    tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{idx:02d}"
    run.font.size = Pt(max(9, d * 26))
    run.font.bold = True
    run.font.name = font
    run.font.color.rgb = _rgb(readable_on(accent))
    return circ


# ─── Overflow-guarded text ────────────────────────────────────────────────────

def _txt(slide, text, l, t, w, h, size=15.0, bold=False, color="#FFFFFF",
         align=PP_ALIGN.LEFT, italic=False, font="Calibri", fit=True,
         min_size: Optional[float] = None, anchor=None):
    """Add a textbox with word-wrap ON, zero internal margins, and the auto-fit
    guard: if the estimated rendered height exceeds the box, the size steps
    down; at the floor the text is word-clipped. Overflow cannot ship.
    `anchor=MSO_ANCHOR.MIDDLE` vertically centers short text in tall boxes so
    cards never read as top-heavy with dead space below."""
    text = str(text or "")
    if fit:
        size = fit_text_size(text, w, h, size, min_pt=min_size, bold=bold)
        while len(text.split()) > 8 and estimate_text_height_in(text, w, size, bold) > h:
            text = clip_words(text, max(8, len(text.split()) - 6))
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        try:
            tf.vertical_anchor = anchor
        except Exception:
            pass
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return tb


def _bullets(slide, items: List[str], l, t, w, h, size=13.0, color="#FFFFFF",
             marker="→", font="Calibri", space_pt=6.0):
    """Bullet list with the same guard: shrink to fit, then drop items rather
    than overflow."""
    items = [str(b) for b in items if str(b).strip()]
    if not items:
        return None

    def total_height(sz: float, lst: List[str]) -> float:
        return sum(estimate_text_height_in(f"{marker}  {b}", w, sz) for b in lst) \
            + (len(lst) - 1) * space_pt / 72.0

    floor = max(9.0, size * 0.7)
    while size > floor and total_height(size, items) > h:
        size -= 0.5
    while len(items) > 1 and total_height(size, items) > h:
        items = items[:-1]

    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(space_pt)
        run = p.add_run()
        run.text = f"{marker}  {b}"
        run.font.size = Pt(size)
        run.font.name = font
        run.font.color.rgb = _rgb(color)
    return tb


def _sentences(text: str, limit: int = 3) -> List[str]:
    parts = [s.strip() for s in re.split(r"(?<=[.!?])\s+", str(text or ""))
             if len(s.strip()) > 12]
    return parts[:limit]


def _rail_content(sec: Section, n: int = 3) -> List[str]:
    """Side-rail / support content that can NEVER be empty: highlights → stat
    contexts → any short-point content → body/notes sentences → callout. An
    empty panel is a design bug; every container must earn its place."""
    items = [h for h in sec.highlights if str(h).strip()]
    if not items:
        items = [st.context for st in sec.stats if st.context]
    if not items:
        items = sec.fallback_highlights()
    if not items:
        items = _sentences(sec.body) or _sentences(sec.speaker_notes)
    if not items and sec.callout:
        items = [sec.callout]
    return [str(i) for i in items[:n]]


# ─── Photos (cover-cropped, never distorted) ──────────────────────────────────

def _photo(slide, section: Section, l, t, w, h, shadow: bool = True) -> bool:
    """Place the prefetched topic image filling the box edge-to-edge with a
    center crop (PowerPoint srcRect), keeping aspect. Returns True if placed."""
    raw = assets.image_bytes(section)
    if not raw:
        return False
    try:
        pic = slide.shapes.add_picture(io.BytesIO(raw), Inches(l), Inches(t),
                                       Inches(w), Inches(h))
        iw, ih = pic.image.size
        if iw and ih:
            img_aspect = iw / ih
            box_aspect = w / h
            if img_aspect > box_aspect:          # too wide → crop sides
                frac = 1 - box_aspect / img_aspect
                pic.crop_left = frac / 2
                pic.crop_right = frac / 2
            elif img_aspect < box_aspect:        # too tall → crop top/bottom
                frac = 1 - img_aspect / box_aspect
                pic.crop_top = frac / 2
                pic.crop_bottom = frac / 2
        if shadow:
            _soft_shadow(pic)
        return True
    except Exception:
        return False


def _scrim(slide, l, t, w, h, color: str, a1: int = 94, a2: int = 42,
           angle: float = 0.0, color2: str = None):
    """Semi-transparent gradient overlay laid over a full-bleed photo so text
    stays readable — the classic editorial 'color scrim' treatment. Alpha per
    stop needs raw OOXML (python-pptx exposes no fill transparency)."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    _grad(shape, color, color2 or color, angle)
    try:
        grad_fill = shape._element.spPr.find(qn('a:gradFill'))
        stops = grad_fill.findall(qn('a:gsLst') + '/' + qn('a:gs'))
        for gs, alpha in zip(stops, (a1, a2)):
            clr = gs.find(qn('a:srgbClr'))
            if clr is not None:
                clr.append(parse_xml(
                    f'<a:alpha {nsdecls("a")} val="{int(alpha * 1000)}"/>'))
    except Exception:
        pass  # opaque scrim still keeps text readable
    return shape


def _hero_photo(d, slide, sec: Section, drama: bool = False) -> int:
    """Full-bleed photo + dominant-color scrim (dense on the left where the
    text sits, lifting toward the right so the photo shows through). `drama`
    (covers/closings) pushes the contrast harder — near-black over the text,
    the photo punching through clean on the right — plus a bottom scrim so
    footer/tagline text always reads. Returns static shape count (0 = no photo)."""
    if not _photo(slide, sec, 0, 0, SLIDE_W, SLIDE_H, shadow=False):
        return 0
    if drama:
        _scrim(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(d.DOM, 0.42), a1=98, a2=28,
               angle=0.0, color2=_shade(d.DOM, 0.62))
        _scrim(slide, 0, 4.9, SLIDE_W, SLIDE_H - 4.9, _shade(d.DOM, 0.38),
               a1=0, a2=82, angle=90.0)
        return 3
    _scrim(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(d.DOM, 0.55), a1=96, a2=45,
           angle=0.0, color2=_shade(d.DOM, 0.75))
    return 2


# ─── Motion: entrance animations, transitions (raw OOXML) ────────────────────
#
# python-pptx has no animation API, so the <p:timing> and <p:transition> trees
# are injected directly. Effects use nodeType="afterEffect" so the slide
# auto-plays on display, staggered for a build-in feel. Schema order matters:
# transition MUST be appended before timing on the slide element.

_ENTRANCES = {
    "fade":       ("fade", 10),
    "wipe_up":    ("wipe(up)", 22),
    "wipe_right": ("wipe(right)", 22),
    "wipe_left":  ("wipe(left)", 22),
    "blinds":     ("blinds(horizontal)", 3),
    "dissolve":   ("dissolve", 9),
    "circle":     ("circle", 5),
    "diamond":    ("diamond", 7),
    "wheel":      ("wheel(1)", 21),
    "plus":       ("plus", 16),
    "randombar":  ("randombar(vertical)", 18),
    "strips":     ("strips(downRight)", 19),
}

# Keynote-safe motion. Keynote (and older PowerPoint) silently DROP effects it
# can't map — blinds/strips/randombar/wheel entrances and cut/zoom/split
# transitions import as "no animation", which reads as a bug. Both pools are
# limited to effects that survive the PowerPoint → Keynote round-trip.
_DECK_ENTRANCE_KEYS = ["fade", "wipe_up", "wipe_right", "wipe_left", "dissolve",
                       "circle", "diamond"]

_TRANSITIONS = [
    '<p:fade/>',
    '<p:fade thruBlk="1"/>',
    '<p:push dir="l"/>',
    '<p:push dir="u"/>',
    '<p:wipe dir="d"/>',
]


def _anim_effect_par(cid: int, spid: int, delay: int, dur: int, filt: str, preset_id: int) -> str:
    """One staggered entrance effect for a single shape (3-level cTn nest)."""
    return (
        f'<p:par><p:cTn id="{cid}" fill="hold">'
        f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
        f'<p:childTnLst><p:par><p:cTn id="{cid + 1}" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst><p:par>'
        f'<p:cTn id="{cid + 2}" presetID="{preset_id}" presetClass="entr" presetSubtype="0" '
        f'fill="hold" nodeType="afterEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'<p:set><p:cBhvr>'
        f'<p:cTn id="{cid + 3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="{filt}"><p:cBhvr>'
        f'<p:cTn id="{cid + 4}" dur="{dur}"/>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    )


def _build_timing(specs: list):
    """specs: list of (spid, delay_ms, dur_ms, filter, preset_id) → <p:timing>."""
    cid = 3
    pars = []
    for spid, delay, dur, filt, preset in specs:
        pars.append(_anim_effect_par(cid, spid, delay, dur, filt, preset))
        cid += 5
    xml = (
        f'<p:timing {nsdecls("p", "a")}><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(pars)}</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    return parse_xml(xml)


def _shuffle_no_adjacent(items: list, rng: random.Random, key=lambda x: x) -> list:
    """Shuffle so no two neighbours share a key — including the wrap-around,
    since the list is cycled across slides (push(l) next to push(u) still
    reads as 'push, push')."""
    items = list(items)
    for _ in range(40):
        rng.shuffle(items)
        n = len(items)
        if n < 2 or all(key(items[i]) != key(items[(i + 1) % n]) for i in range(n)):
            break
    return items


def _set_transition(slide, transition_xml: str):
    """Inject <p:transition>. MUST run before _animate to keep schema order."""
    xml = f'<p:transition {nsdecls("p")} spd="med">{transition_xml}</p:transition>'
    slide._element.append(parse_xml(xml))


def _animate(slide, filt_key: str, max_shapes: int = 9, step: int = 170, dur: int = 460,
             skip: int = 1):
    """Auto-playing staggered entrance. `skip` = leading static shapes
    (background rect, or full-bleed photo + scrim) that must not animate."""
    shapes = list(slide.shapes)
    if len(shapes) <= skip:
        return
    filt, preset = _ENTRANCES.get(filt_key, _ENTRANCES["fade"])
    specs = [(sh.shape_id, 100 + i * step, dur, filt, preset)
             for i, sh in enumerate(shapes[skip:skip + max_shapes])]
    if specs:
        slide._element.append(_build_timing(specs))


def _nav_controls(slide, prs, idx: int, total: int, theme: Theme):
    """Clickable prev / home / next dots that jump between slides."""
    targets = []
    if idx > 0:
        targets.append(("‹", idx - 1))
    targets.append(("●", 0))
    if idx < total - 1:
        targets.append(("›", idx + 1))
    bx = SLIDE_W - 0.05 - len(targets) * 0.42
    on_acc = readable_on(theme.accent)
    for label, tgt in targets:
        btn = _oval(slide, bx, 7.02, 0.34, theme.accent)
        tf = btn.text_frame
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = _rgb(on_acc)
        try:
            btn.click_action.target_slide = prs.slides[tgt]
        except Exception:
            pass
        bx += 0.42


# ─── Native editable chart ────────────────────────────────────────────────────

def _native_chart(slide, section: Section, l, t, w, h, theme: Theme, on_dark: bool = True):
    chart_block = section.chart
    if not chart_block or len(chart_block.points) < 2:
        return None
    cd = CategoryChartData()
    cd.categories = [str(lbl)[:14] for lbl, _ in chart_block.points]
    cd.add_series("Series 1", [v for _, v in chart_block.points])
    kind = "line" if ("trend" in (section.heading + chart_block.kind).lower()
                      or "growth" in (section.heading + chart_block.kind).lower()
                      or chart_block.kind == "line") else "column"
    ctype = XL_CHART_TYPE.LINE_MARKERS if kind == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    frame = slide.shapes.add_chart(ctype, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = frame.chart
    label_color = _rgb("#FFFFFF" if on_dark else theme.ink)
    tick_color = _rgb(theme.tint if on_dark else theme.body_ink)
    data_font = theme.pptx_fonts["data"]
    try:
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.gap_width = 60
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(10)
        plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = label_color
        plot.data_labels.font.name = data_font
        series = plot.series[0]
        if kind == "line":
            series.format.line.color.rgb = _rgb(theme.accent)
            series.format.line.width = Pt(2.75)
            try:
                series.smooth = True   # curved vector line, not a polyline
            except Exception:
                pass
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(theme.accent)
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.color.rgb = tick_color
        cat_axis.tick_labels.font.name = data_font
        cat_axis.format.line.color.rgb = tick_color
        val_axis = chart.value_axis
        val_axis.visible = False
        val_axis.has_major_gridlines = False
    except Exception:
        pass  # styling is best-effort; the chart itself still renders
    return frame


# ─── Layout planning ──────────────────────────────────────────────────────────

def plan_pptx_layouts(sections: List[Section], seed: int) -> None:
    """Content-compatible layout per slide, shuffled per deck, consecutive
    repeats repaired — the same content never renders the same way twice."""
    rng = random.Random(seed)

    def compatible(sec: Section) -> list:
        opts = []
        if sec.chart:
            opts.append("chart")
        if len(sec.steps) >= 3:
            opts.append("timeline")
        if len(sec.tiles) >= 3:
            opts.append("cards")
        if sec.left and sec.right:
            opts.append("two_col")
        if len(sec.stats) >= 2:
            opts.append("stat")
        if sec.callout or sec.attribution:
            opts.append("quote")
        if sec.body and (sec.fallback_highlights() or sec.image_data):
            opts.append("image_text")
        if not opts:
            opts = ["stat" if sec.stats else "image_text"]
        return opts

    options, assigned = [], []
    n = len(sections)
    for i, sec in enumerate(sections):
        if i == 0 or sec.kind == "cover":
            options.append(["title"]); assigned.append("title"); continue
        if i == n - 1 or sec.kind == "closing":
            options.append(["closing"]); assigned.append("closing"); continue
        opts = compatible(sec)
        rng.shuffle(opts)
        prev = assigned[-1] if assigned else ""
        # prefer a layout unused so far in the deck, then any non-repeat
        choice = next((o for o in opts if o != prev and o not in assigned),
                      next((o for o in opts if o != prev), opts[0]))
        options.append(opts)
        assigned.append(choice)

    for sec, layout in zip(sections, no_consecutive_repeat(assigned, options, rng)):
        sec.layout = layout


# ─── Slide builders ───────────────────────────────────────────────────────────

class _Deck:
    """Carries theme-resolved colors/fonts + per-deck randomness for builders."""

    def __init__(self, prs: Presentation, spec: DocumentSpec, theme: Theme):
        self.prs = prs
        self.spec = spec
        self.theme = theme
        self.total = len(spec.sections)
        t = theme
        self.F = theme.pptx_fonts                      # display/heading/body/data
        self.S = theme.scale
        self.DOM, self.SEC, self.ACC = t.dominant, t.secondary, t.accent
        self.acc_dark = t.accent_readable_on_dark()    # accent text on dark bg
        self.acc_light = t.accent_readable_on_light()  # accent text on light bg
        self.INK, self.BODY = t.ink, t.body_ink
        self.TINT, self.WASH = t.tint, t.wash
        self.MUTED_DARK = _shade(t.support, 1.0)       # muted text on dark
        self.dark_deck = t.dark_deck

    # backgrounds ---------------------------------------------------------

    def dark_bg(self, slide):
        _grad_rect(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(self.DOM, 0.66), self.DOM, 125)

    def light_bg(self, slide):
        _grad_rect(slide, 0, 0, SLIDE_W, SLIDE_H, "#FFFFFF", self.WASH, 125)

    def content_bg(self, slide) -> bool:
        """Returns True if the content surface is dark."""
        if self.dark_deck:
            self.dark_bg(slide)
            return True
        self.light_bg(slide)
        return False

    def photo_band(self, slide, sec: Section) -> int:
        """Full-width photo strip with scrim behind the eyebrow/title zone —
        turns any text-heavy layout into an editorial header. Returns the
        number of static shapes added (0 when the section has no image)."""
        if not _photo(slide, sec, 0, 0, SLIDE_W, 2.02, shadow=False):
            return 0
        _scrim(slide, 0, 0, SLIDE_W, 2.02, _shade(self.DOM, 0.5), a1=95, a2=62,
               angle=0.0, color2=_shade(self.DOM, 0.68))
        return 2

    # furniture -----------------------------------------------------------

    def footer(self, slide, idx: int, dark: bool, right_edge: float = None):
        muted = "#FFFFFF" if dark else self.BODY
        acc = self.acc_dark if dark else self.acc_light
        right_edge = SLIDE_W - MARGIN if right_edge is None else right_edge
        title = clip_chars(self.spec.title.upper(), 70)
        _txt(slide, title, MARGIN, 7.05, min(9.0, right_edge - MARGIN - 1.7), 0.3,
             size=8, color=muted, font=self.F["data"], fit=False)
        _txt(slide, f"{idx + 1:02d} / {self.total:02d}", right_edge - 1.5, 7.05,
             1.5, 0.3, size=8, color=acc, align=PP_ALIGN.RIGHT, font=self.F["data"], fit=False)

    def eyebrow(self, slide, idx: int, text: str, dark: bool):
        acc = self.acc_dark if dark else self.acc_light
        _num_circle(slide, MARGIN, 0.48, 0.34, idx + 1, self.ACC, self.F["data"])
        _txt(slide, clip_chars(text, 44).upper(), MARGIN + 0.5, 0.53, 11, 0.32,
             size=9.5, bold=True, color=acc, font=self.F["data"], fit=False)

    def slide_title(self, slide, text: str, dark: bool, width: float = 12.1):
        color = "#FFFFFF" if dark else self.INK
        _txt(slide, text, MARGIN - 0.02, 1.0, width, 1.15, size=self.S.slide_title,
             bold=True, color=color, font=self.F["heading"], min_size=24)


def _build_title(d: _Deck, slide, sec: Section, idx: int):
    # Full-bleed topic photo behind a color scrim when available — the single
    # biggest "designed, not generated" signal a cover can send.
    has_photo = _hero_photo(d, slide, sec, drama=True)
    if has_photo:
        static = has_photo
        static += _dot_grid(slide, 11.05, 5.55, 6, 4, d.acc_dark)
        _outline_ring(slide, 10.15, 1.05, 2.9, d.acc_dark, weight=2.0)
        static += 1
    else:
        _grad_rect(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(d.DOM, 0.6), d.DOM, 120)
        static = 1 + _dot_grid(slide, 0.62, 5.9, 8, 3, _shade(d.SEC, 1.35))
        ring = slide.shapes.add_shape(9, Inches(9.4), Inches(-1.8), Inches(6.2), Inches(6.2))
        _grad(ring, _shade(d.SEC, 1.12), d.SEC, 120); _soft_shadow(ring, blur_pt=24, dist_pt=0, alpha_pct=22)
        ring2 = slide.shapes.add_shape(9, Inches(11.1), Inches(4.7), Inches(3.4), Inches(3.4))
        _grad(ring2, d.ACC, _shade(d.ACC, 0.7), 120); _soft_shadow(ring2, blur_pt=18, dist_pt=0, alpha_pct=26)

    tagline = sec.tagline
    if d.theme.aesthetic_label and d.theme.aesthetic_label not in tagline:
        tagline = f"{d.theme.aesthetic_label}  |  {tagline}".strip(" |")
    if d.theme.aesthetic_label:
        # kicker chip — small accent-filled label above the headline
        label = clip_chars(d.theme.aesthetic_label.upper(), 30)
        chip_w = min(4.6, 0.42 + len(label) * 0.085)
        chip = _round_card(slide, MARGIN, 0.92, chip_w, 0.42, color=d.ACC,
                           radius=0.5, shadow=False)
        ctf = chip.text_frame
        ctf.word_wrap = False
        ctf.paragraphs[0].text = label
        r0 = ctf.paragraphs[0].runs[0]
        r0.font.size = Pt(10); r0.font.bold = True
        r0.font.color.rgb = _rgb(readable_on(d.ACC)); r0.font.name = d.F["data"]
    _txt(slide, sec.heading or d.spec.title, MARGIN - 0.02, 1.65, 9.6, 3.0,
         size=(d.S.cover_title + 8 if has_photo else d.S.cover_title), bold=True,
         color="#FFFFFF", font=d.F["display"], min_size=28)
    subtitle = sec.subtitle or d.spec.subtitle
    if subtitle:
        _txt(slide, subtitle, MARGIN, 4.9, 9.2, 1.1, size=20, color=d.TINT, font=d.F["body"])
    if tagline:
        _txt(slide, tagline, MARGIN, 6.65, 10.5, 0.5, size=11.5, color=d.MUTED_DARK, font=d.F["body"])
    _txt(slide, f"01 / {d.total:02d}", SLIDE_W - 1.5 - MARGIN, 6.7, 1.5, 0.4, size=11,
         color=d.acc_dark, align=PP_ALIGN.RIGHT, font=d.F["data"], fit=False)
    return static


def _build_closing(d: _Deck, slide, sec: Section, idx: int):
    static = _hero_photo(d, slide, sec, drama=True)
    if static:
        static += _dot_grid(slide, 11.05, 0.62, 6, 4, d.acc_dark)
        _outline_ring(slide, 10.5, 4.6, 2.4, d.acc_dark, weight=2.0)
        static += 1
    else:
        static = 1
        _grad_rect(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(d.DOM, 0.6), d.DOM, 120)
        static += _dot_grid(slide, 0.62, 6.15, 8, 2, _shade(d.SEC, 1.35))
        ring = slide.shapes.add_shape(9, Inches(8.2), Inches(0.4), Inches(6.2), Inches(6.2))
        _grad(ring, _shade(d.SEC, 1.1), d.SEC, 120); _soft_shadow(ring, blur_pt=24, dist_pt=0, alpha_pct=22)
    _txt(slide, "IN CLOSING", MARGIN, 0.7, 5, 0.4, size=10, bold=True,
         color=d.acc_dark, font=d.F["data"], fit=False)
    _txt(slide, sec.heading or "Thank you", MARGIN - 0.02, 1.5, 9.4, 2.5, size=40,
         bold=True, color="#FFFFFF", font=d.F["display"], min_size=26)
    cta = sec.cta or sec.callout or "Thank you."
    cta_box = _round_card(slide, MARGIN, 4.35, 7.6, 1.05,
                          grad=(d.ACC, _shade(d.ACC, 0.82), 90), radius=0.16)
    ctf = cta_box.text_frame
    ctf.word_wrap = True
    ctf.paragraphs[0].text = clip_words(cta, 18)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r0 = ctf.paragraphs[0].runs[0]
    r0.font.size = Pt(fit_text_size(cta, 7.0, 0.9, 17, min_pt=12, bold=True))
    r0.font.bold = True
    r0.font.color.rgb = _rgb(readable_on(d.ACC))
    r0.font.name = d.F["body"]
    if sec.tagline:
        _txt(slide, sec.tagline, MARGIN, 5.65, 9, 0.6, size=13.5, color=d.TINT,
             italic=True, font=d.F["body"])
    d.footer(slide, idx, dark=True)
    return static


def _build_stat(d: _Deck, slide, sec: Section, idx: int):
    d.dark_bg(slide)          # stat slides always sit on the dominant color
    dark = True
    static = 1 + d.photo_band(slide, sec)
    if static == 1:           # no photo band → vector dot accent instead
        static += _dot_grid(slide, 11.35, 0.52, 5, 3, _shade(d.SEC, 1.3))
    d.eyebrow(slide, idx, sec.eyebrow or "KEY METRICS", dark)
    d.slide_title(slide, sec.heading, dark)
    stats = sec.stats[:3]
    card_w = (SLIDE_W - 2 * MARGIN - 2 * GAP) / 3
    for ci, st in enumerate(stats):
        cx = MARGIN + ci * (card_w + GAP)
        _round_card(slide, cx, 2.45, card_w, 4.0, grad=(_shade(d.SEC, 1.12), d.SEC, 120), radius=0.05)
        _num_circle(slide, cx + 0.32, 2.75, 0.4, ci + 1, d.ACC, d.F["data"])
        _txt(slide, st.value, cx + 0.25, 3.3, card_w - 0.5, 1.4, size=d.S.stat_card,
             bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER, font=d.F["data"], min_size=24)
        _txt(slide, st.label.upper(), cx + 0.25, 4.8, card_w - 0.5, 0.5, size=12.5,
             bold=True, color=d.acc_dark, align=PP_ALIGN.CENTER, font=d.F["heading"], min_size=9)
        pct = _pct_of(st.value)
        if pct is not None:   # percentage stat → vector gauge bar
            _progress_bar(slide, cx + 0.32, 5.38, card_w - 0.64, pct,
                          d.ACC, _shade(d.SEC, 1.45))
            _txt(slide, st.context, cx + 0.3, 5.68, card_w - 0.6, 0.7, size=10,
                 color=d.TINT, align=PP_ALIGN.LEFT, font=d.F["body"])
        else:
            _txt(slide, st.context, cx + 0.3, 5.35, card_w - 0.6, 0.95, size=10.5,
                 color=d.TINT, align=PP_ALIGN.LEFT, font=d.F["body"])
    d.footer(slide, idx, dark)
    return static


def _build_two_col(d: _Deck, slide, sec: Section, idx: int):
    dark = d.content_bg(slide)
    band = d.photo_band(slide, sec)
    head_dark = dark or bool(band)
    d.eyebrow(slide, idx, sec.eyebrow or "COMPARISON", head_dark)
    d.slide_title(slide, sec.heading, head_dark)
    col_w = (SLIDE_W - 2 * MARGIN - GAP) / 2
    left_x, right_x = MARGIN, MARGIN + col_w + GAP
    # left: light card; right: deep card — contrast without any edge stripes
    if dark:
        _round_card(slide, left_x, 2.15, col_w, 4.6, grad=(_shade(d.SEC, 1.16), _shade(d.SEC, 0.92), 120), radius=0.045)
        left_text = "#FFFFFF"
    else:
        _round_card(slide, left_x, 2.15, col_w, 4.6, grad=("#FFFFFF", d.TINT, 120), radius=0.045)
        left_text = d.BODY
    _round_card(slide, right_x, 2.15, col_w, 4.6, grad=(_shade(d.SEC, 1.08), d.SEC, 120), radius=0.045)
    lt = (sec.left.title if sec.left else "") or "Perspective A"
    rt = (sec.right.title if sec.right else "") or "Perspective B"
    acc_on_card = d.acc_dark if dark else d.acc_light
    _txt(slide, lt.upper(), left_x + 0.35, 2.5, col_w - 0.7, 0.4, size=11, bold=True,
         color=acc_on_card if not dark else d.acc_dark, font=d.F["data"], fit=False)
    _txt(slide, rt.upper(), right_x + 0.35, 2.5, col_w - 0.7, 0.4, size=11, bold=True,
         color=d.acc_dark, font=d.F["data"], fit=False)
    _bullets(slide, (sec.left.points if sec.left else sec.highlights)[:4],
             left_x + 0.35, 3.05, col_w - 0.7, 3.5, size=13.5, color=left_text,
             marker="→", font=d.F["body"])
    _bullets(slide, (sec.right.points if sec.right else sec.highlights)[:4],
             right_x + 0.35, 3.05, col_w - 0.7, 3.5, size=13.5, color="#FFFFFF",
             marker="→", font=d.F["body"])
    d.footer(slide, idx, dark)
    return 1 + band


def _build_chart(d: _Deck, slide, sec: Section, idx: int):
    d.dark_bg(slide)
    band = d.photo_band(slide, sec)
    d.eyebrow(slide, idx, sec.eyebrow or "DATA", True)
    d.slide_title(slide, sec.heading, True)
    panel = "#0F1218"
    _round_card(slide, MARGIN - 0.05, 2.35, 8.35, 4.35, grad=(_shade(panel, 1.5), panel, 120), radius=0.04)
    frame = _native_chart(slide, sec, 0.85, 2.7, 7.75, 3.7, d.theme, on_dark=True)
    if frame is None:
        for ci, st in enumerate(sec.stats[:3]):
            cx = [0.85, 3.55, 6.25][ci]
            _txt(slide, st.value, cx, 3.4, 2.5, 1.2, size=34, bold=True,
                 color="#FFFFFF", font=d.F["data"], min_size=20)
            _txt(slide, st.label, cx, 4.6, 2.5, 0.6, size=11, bold=True,
                 color=d.acc_dark, font=d.F["heading"])
    _round_card(slide, 9.15, 2.35, 3.6, 4.35, grad=(_shade(d.SEC, 1.1), d.SEC, 120), radius=0.05)
    _num_circle(slide, 9.4, 2.65, 0.38, idx + 1, d.ACC, d.F["data"])
    _txt(slide, "WHAT IT MEANS", 9.9, 2.72, 2.7, 0.35, size=10, bold=True,
         color=d.acc_dark, font=d.F["data"], fit=False)
    _bullets(slide, _rail_content(sec, 4), 9.4, 3.3, 3.1, 3.2, size=11.5,
             color="#FFFFFF", marker="◆", font=d.F["body"])
    d.footer(slide, idx, True)
    return 1 + band


def _build_quote(d: _Deck, slide, sec: Section, idx: int):
    # Photo-backed quotes get a near-opaque scrim: the image adds atmosphere,
    # the words stay the hero.
    static = 1
    if _photo(slide, sec, 0, 0, SLIDE_W, SLIDE_H, shadow=False):
        static = 2
        _scrim(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(d.DOM, 0.5), a1=93, a2=82,
               angle=90.0, color2=_shade(d.DOM, 0.66))
    else:
        _grad_rect(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(d.DOM, 0.62), d.DOM, 125)
    _txt(slide, "“", 0.45, -0.35, 3, 2.6, size=150, color=_shade(d.SEC, 1.2), bold=True,
         font=d.F["display"], fit=False)
    quote = sec.callout or sec.body
    _txt(slide, quote, 1.0, 1.75, 8.9, 3.6, size=d.S.quote, color="#FFFFFF",
         italic=True, font=d.F["display"], min_size=16)
    if sec.attribution:
        _txt(slide, sec.attribution, 1.05, 5.6, 8.5, 0.5, size=13, color=d.acc_dark,
             bold=True, font=d.F["data"])
    hi = [x for x in (sec.highlights or sec.statements) if str(x).strip()] or \
        _rail_content(sec, 3)
    hi = [h for h in hi if str(h).strip() != str(quote).strip()]
    for hj, h in enumerate(hi[:3]):
        hy = 1.35 + hj * 1.75
        _round_card(slide, 10.2, hy, 2.5, 1.5, grad=(_shade(d.SEC, 1.1), d.SEC, 120), radius=0.08)
        _num_circle(slide, 10.4, hy + 0.18, 0.3, hj + 1, d.ACC, d.F["data"])
        _txt(slide, str(h), 10.4, hy + 0.58, 2.1, 0.82, size=10.5, color="#FFFFFF",
             font=d.F["body"])
    d.footer(slide, idx, True)
    return static


def _build_image_text(d: _Deck, slide, sec: Section, idx: int):
    dark = d.content_bg(slide)
    # Media hierarchy: embedded video (16:9 panel, poster = section photo) >
    # editorial full-height photo split > geometric data panel.
    video = assets.video_bytes(sec)
    vid_x, vid_y, vid_w = 6.6, 2.3, SLIDE_W - 6.6 - MARGIN
    vid_h = vid_w * 9.0 / 16.0
    split_x = 8.15
    has_photo = False
    if not video:
        has_photo = _photo(slide, sec, split_x, 0, SLIDE_W - split_x, SLIDE_H, shadow=False)
    text_right = (vid_x - 0.4) if video else (split_x - 0.45 if has_photo else 8.5)
    d.eyebrow(slide, idx, sec.eyebrow or "OVERVIEW", dark)
    d.slide_title(slide, sec.heading, dark, width=text_right - MARGIN)
    body_color = "#FFFFFF" if dark else d.BODY
    text_w = text_right - MARGIN
    # structured lower half: body up top, then numbered insight rows — a short
    # paragraph must never leave half a slide of dead space
    chip_text = sec.callout or (sec.fallback_highlights() or [""])[0]
    rows = [r for r in _rail_content(sec, 4)
            if r not in (sec.body or "") and r.strip() != str(chip_text).strip()][:3] \
        if (has_photo or video) else []   # no-media branch shows these in its side panel
    body_h = 2.35 if rows else 4.4
    _txt(slide, sec.body, MARGIN, 2.3, text_w, body_h, size=d.S.body,
         color=body_color, font=d.F["body"], min_size=11)
    if rows:
        _rect(slide, MARGIN, 4.85, 1.1, 0.03, d.ACC)   # accent divider
        row_txt = "#FFFFFF" if dark else d.BODY
        for rj, r in enumerate(rows):
            ry = 5.05 + rj * 0.62
            _num_circle(slide, MARGIN, ry, 0.32, rj + 1, d.ACC, d.F["data"])
            _txt(slide, str(r), MARGIN + 0.5, ry + 0.015, text_w - 0.5, 0.56,
                 size=11.5, color=row_txt, font=d.F["body"], min_size=9,
                 anchor=MSO_ANCHOR.MIDDLE)
    if video:
        placed = False
        try:
            poster = assets.image_bytes(sec)
            movie = slide.shapes.add_movie(
                io.BytesIO(video), Inches(vid_x), Inches(vid_y),
                Inches(vid_w), Inches(vid_h),
                poster_frame_image=io.BytesIO(poster) if poster else None,
                mime_type="video/mp4")
            _soft_shadow(movie)
            placed = True
        except Exception:
            # video part failed to embed — the section photo fills the slot
            placed = _photo(slide, sec, vid_x, vid_y, vid_w, vid_h)
        if placed and str(chip_text).strip():
            chip = _round_card(slide, vid_x, vid_y + vid_h + 0.25, vid_w, 0.9,
                               color=_shade(d.DOM, 0.85), radius=0.12, shadow=True)
            ctf = chip.text_frame
            ctf.word_wrap = True
            try:
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            ctf.paragraphs[0].text = clip_words(str(chip_text), 16)
            r0 = ctf.paragraphs[0].runs[0]
            r0.font.size = Pt(10.5)
            r0.font.color.rgb = _rgb("#FFFFFF")
            r0.font.name = d.F["body"]
        d.footer(slide, idx, dark)
    elif has_photo:
        if str(chip_text).strip():
            chip = _round_card(slide, split_x + 0.35, SLIDE_H - 1.5, SLIDE_W - split_x - 0.9,
                               1.0, color=_shade(d.DOM, 0.85), radius=0.12, shadow=True)
            ctf = chip.text_frame
            ctf.word_wrap = True
            try:
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            ctf.paragraphs[0].text = clip_words(str(chip_text), 14)
            r0 = ctf.paragraphs[0].runs[0]
            r0.font.size = Pt(10.5)
            r0.font.color.rgb = _rgb("#FFFFFF")
            r0.font.name = d.F["body"]
        d.footer(slide, idx, dark, right_edge=text_right)
    else:
        vis_x, vis_y, vis_w, vis_h = 8.9, 2.15, SLIDE_W - 8.9 - MARGIN, 4.6
        _round_card(slide, vis_x, vis_y, vis_w, vis_h, grad=(_shade(d.SEC, 1.1), d.SEC, 120), radius=0.05)
        _num_circle(slide, vis_x + 0.3, vis_y + 0.3, 0.38, idx + 1, d.ACC, d.F["data"])
        _txt(slide, "KEY DATA", vis_x + 0.82, vis_y + 0.37, vis_w - 1.1, 0.35, size=10,
             bold=True, color=d.acc_dark, font=d.F["data"], fit=False)
        _bullets(slide, sec.fallback_highlights()[:4], vis_x + 0.3, vis_y + 0.95, vis_w - 0.6,
                 vis_h - 1.25, size=12, color="#FFFFFF", marker="◆", font=d.F["body"])
        d.footer(slide, idx, dark)


def _build_cards(d: _Deck, slide, sec: Section, idx: int):
    """Bento grid: one featured accent tile + varied-size neighbours instead of
    a monotonous 2x2 — the asymmetry is what reads as designed."""
    dark = d.content_bg(slide)
    band = d.photo_band(slide, sec)
    head_dark = dark or bool(band)
    d.eyebrow(slide, idx, sec.eyebrow or "TAKEAWAYS", head_dark)
    d.slide_title(slide, sec.heading, head_dark)
    items = [(t.label, t.value) for t in sec.tiles[:4]] or \
            [("", h) for h in sec.fallback_highlights()[:4]]

    top, total_h = 2.3, 4.35
    total_w = SLIDE_W - 2 * MARGIN
    feat_w = 4.9
    side_w = total_w - feat_w - GAP
    half_h = (total_h - GAP) / 2
    half_w = (side_w - GAP) / 2
    mirror = (d.theme.seed + idx) % 2 == 1
    feat_x = MARGIN + (total_w - feat_w if mirror else 0)
    side_x = MARGIN + (0 if mirror else feat_w + GAP)

    # cell geometry per item count: featured tall + wide + squares
    if len(items) >= 4:
        cells = [(feat_x, top, feat_w, total_h),
                 (side_x, top, side_w, half_h),
                 (side_x, top + half_h + GAP, half_w, half_h),
                 (side_x + half_w + GAP, top + half_h + GAP, half_w, half_h)]
    elif len(items) == 3:
        cells = [(feat_x, top, feat_w, total_h),
                 (side_x, top, side_w, half_h),
                 (side_x, top + half_h + GAP, side_w, half_h)]
    else:
        cells = [(MARGIN + i * (total_w + GAP) / 2, top,
                  (total_w - GAP) / 2, total_h) for i in range(len(items))]

    for i, ((label, value), (cx, cy, cw, ch)) in enumerate(zip(items, cells)):
        featured = i == 0 and len(items) >= 3
        if featured:
            _round_card(slide, cx, cy, cw, ch, grad=(d.ACC, _shade(d.ACC, 0.78), 120),
                        radius=0.07)
            txt_color = readable_on(d.ACC)
            circ_color, label_color = d.DOM, txt_color
        elif dark:
            _round_card(slide, cx, cy, cw, ch, grad=(_shade(d.SEC, 1.12), d.SEC, 120), radius=0.07)
            txt_color, circ_color, label_color = "#FFFFFF", d.ACC, d.acc_dark
        else:
            _round_card(slide, cx, cy, cw, ch, grad=("#FFFFFF", d.TINT, 120), radius=0.07)
            txt_color, circ_color, label_color = d.BODY, d.ACC, d.acc_light
        _num_circle(slide, cx + 0.3, cy + 0.3, 0.42, i + 1, circ_color, d.F["data"])
        if label and not label.isdigit():
            _txt(slide, label.upper(), cx + 0.92, cy + 0.38, cw - 1.2, 0.35,
                 size=10.5, bold=True, color=label_color, font=d.F["data"], fit=False)
        body_size = 17.0 if featured else 12.5
        _txt(slide, value, cx + 0.35 if featured else cx + 0.92,
             cy + 0.95 if featured else cy + 0.82,
             cw - 0.7 if featured else cw - 1.2,
             ch - 1.25 if featured else ch - 1.05,
             size=body_size, bold=featured, color=txt_color, font=d.F["body"],
             min_size=10, anchor=MSO_ANCHOR.MIDDLE)
    d.footer(slide, idx, dark)
    return 1 + band


def _build_timeline(d: _Deck, slide, sec: Section, idx: int):
    dark = d.content_bg(slide)
    band = d.photo_band(slide, sec)
    head_dark = dark or bool(band)
    d.eyebrow(slide, idx, sec.eyebrow or "TIMELINE", head_dark)
    d.slide_title(slide, sec.heading, head_dark)
    steps = sec.steps[:4]
    # connector spine behind the cards — visible in the gaps, reads as a
    # timeline instead of four floating boxes
    _rect(slide, MARGIN + 0.2, 2.87, SLIDE_W - 2 * (MARGIN + 0.2), 0.035,
          d.ACC if dark else _shade(d.ACC, 1.0))
    card_w = (SLIDE_W - 2 * MARGIN - (len(steps) - 1) * GAP) / max(1, len(steps))
    for i, step in enumerate(steps):
        cx = MARGIN + i * (card_w + GAP)
        if dark:
            _round_card(slide, cx, 2.4, card_w, 4.3, grad=(_shade(d.SEC, 1.12), d.SEC, 120), radius=0.06)
            txt_color, date_color = "#FFFFFF", d.acc_dark
        else:
            _round_card(slide, cx, 2.4, card_w, 4.3, grad=("#FFFFFF", d.TINT, 120), radius=0.06)
            txt_color, date_color = d.BODY, d.acc_light
        _num_circle(slide, cx + 0.3, 2.7, 0.42, i + 1, d.ACC, d.F["data"])
        _txt(slide, clip_chars(step.date or f"{i + 1:02d}", 14), cx + 0.3, 3.35,
             card_w - 0.6, 0.45, size=15, bold=True, color=date_color,
             font=d.F["data"], min_size=11)
        _txt(slide, clip_chars(step.title, 30).upper(), cx + 0.3, 3.85, card_w - 0.6,
             0.65, size=11.5, bold=True, color="#FFFFFF" if dark else d.INK,
             font=d.F["heading"], min_size=9)
        _txt(slide, step.description, cx + 0.3, 4.55, card_w - 0.6, 1.95, size=11,
             color=txt_color, font=d.F["body"], min_size=9)
    d.footer(slide, idx, dark)
    return 1 + band


_BUILDERS = {
    "title": _build_title,
    "timeline": _build_timeline,
    "closing": _build_closing,
    "stat": _build_stat,
    "two_col": _build_two_col,
    "chart": _build_chart,
    "quote": _build_quote,
    "image_text": _build_image_text,
    "cards": _build_cards,
}


# ─── Entry point ──────────────────────────────────────────────────────────────

def _register_notes_master(prs) -> None:
    """python-pptx creates the notes-master part and relationship when the
    first notes slide is added, but never lists it in presentation.xml.
    PowerPoint ignores the orphaned relationship; Keynote (macOS and iOS)
    rejects the whole file as "invalid format"."""
    pres = prs.element
    if pres.find(qn("p:notesMasterIdLst")) is not None:
        return
    rid = next((rel.rId for rel in prs.part.rels.values()
                if rel.reltype == RT.NOTES_MASTER), None)
    if rid is None:
        return
    lst = parse_xml(
        f'<p:notesMasterIdLst {nsdecls("p", "r")}>'
        f'<p:notesMasterId r:id="{rid}"/></p:notesMasterIdLst>'
    )
    # schema order: sldMasterIdLst, notesMasterIdLst, …, sldIdLst, sldSz
    sld_id_lst = pres.find(qn("p:sldIdLst"))
    if sld_id_lst is not None:
        sld_id_lst.addprevious(lst)
    else:
        pres.insert(1, lst)


def render_pptx(spec: DocumentSpec, theme: Theme, seed: Optional[int] = None) -> bytes:
    prs = Presentation()
    # exact 16:9; sldSz must be a whole multiple of 12700 EMU (1pt), which
    # Inches(13.333) misses by 305 EMU
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    deck_seed = theme.seed if seed is None else seed
    plan_pptx_layouts(spec.sections, deck_seed)

    # Motion identity — freshly randomized every generation AND varied per
    # slide: both pools are shuffled with the deck seed, then cycled, so two
    # runs get different effect sequences and consecutive slides never share
    # the same entrance or transition.
    motion_rng = random.Random(deck_seed)
    entrance_keys = _shuffle_no_adjacent(
        _DECK_ENTRANCE_KEYS, motion_rng, key=lambda k: k.split('_')[0])
    transitions = _shuffle_no_adjacent(
        _TRANSITIONS, motion_rng, key=lambda t: t.split()[0].strip('<>/'))
    anim_step = motion_rng.choice((140, 170, 200))   # stagger rhythm per deck
    anim_dur = motion_rng.choice((420, 460, 520))    # entrance speed per deck

    d = _Deck(prs, spec, theme)
    for si, sec in enumerate(spec.sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        builder = _BUILDERS.get(sec.layout, _build_image_text)
        # builders return the count of leading static shapes (bg / photo+scrim)
        static = builder(d, slide, sec, si) or 1

        # speaker notes on every slide
        notes = sec.speaker_notes or (
            f"{sec.heading}. " + clip_words(sec.body or sec.callout or spec.subtitle, 40))
        slide.notes_slide.notes_text_frame.text = notes

        # motion: transition MUST precede timing in the slide XML
        _set_transition(slide, transitions[si % len(transitions)])
        _animate(slide, entrance_keys[si % len(entrance_keys)],
                 step=anim_step, dur=anim_dur, skip=static)

    # clickable navigation — added after all slides exist so targets resolve
    all_slides = list(prs.slides)
    for ni, slide in enumerate(all_slides):
        _nav_controls(slide, prs, ni, len(all_slides), theme)

    # professional metadata
    core = prs.core_properties
    core.title = spec.title
    core.author = spec.author or "AI Document Studio"
    core.subject = spec.subtitle or spec.title
    core.comments = f"Theme: {theme.name} · {theme.aesthetic_label}"
    core.created = datetime.now(timezone.utc).replace(tzinfo=None)

    _register_notes_master(prs)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
